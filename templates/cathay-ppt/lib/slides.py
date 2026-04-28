"""
Cathay PPT Template v3 — Slide Structure
==========================================
Uses 阿维塔 template Layout [4] with left red vertical line (5mm).
Title: dark red (#800000) bold, left-aligned.
Content zone: CL=14mm (after red line), CT=31mm to CB=181mm.
"""

from pptx.util import Mm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from constants import (
    CL, CT, CW, CH, CB,
    CATHAY_RED, CATHAY_GOLD, CATHAY_DARK_RED, CATHAY_WHITE, CATHAY_BLACK, CATHAY_GREY, CATHAY_DARK_GREY,
    SOURCE_Y_MM, SOURCE_BOX_HEIGHT_MM, SOURCE_FONT_PT,
    TITLE_FONT_PT, SUBTITLE_FONT_PT, BODY_FONT_PT, SMALL_FONT_PT, CAPTION_FONT_PT,
    RED_LINE_WIDTH, GAP_XS, GAP_SM,
)
from fonts import set_run_font, add_mixed_text
from text_layout import setup_text_frame, format_paragraph
from merge import _clean_shape


# ============================================================================
# 1. COVER SLIDE (Layout [0])
# ============================================================================

def create_cover_slide(prs, fund_name="Cathay Smart Energy Fund",
                       company_name="Company Name", subtitle="Investment Memo",
                       date_text=""):
    """Create Layout [0] cover slide."""
    cover = prs.slides.add_slide(prs.slide_layouts[0])

    # Fund name — top-left
    txBox = cover.shapes.add_textbox(Mm(CL), Mm(44), Mm(200), Mm(15))
    tf = txBox.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_font(run, fund_name, size_pt=TITLE_FONT_PT, color_rgb=CATHAY_RED)

    # Company name
    txBox2 = cover.shapes.add_textbox(Mm(CL), Mm(76), Mm(216), Mm(35))
    tf2 = txBox2.text_frame
    setup_text_frame(tf2)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    set_run_font(run2, company_name, size_pt=28, bold=True, color_rgb=CATHAY_BLACK)
    # Subtitle
    p2b = tf2.add_paragraph()
    p2b.alignment = PP_ALIGN.LEFT
    run2b = p2b.add_run()
    set_run_font(run2b, subtitle, size_pt=18, color_rgb=CATHAY_GREY)

    # Date
    if date_text:
        txBox3 = cover.shapes.add_textbox(Mm(CL), Mm(127), Mm(76), Mm(10))
        tf3 = txBox3.text_frame
        setup_text_frame(tf3)
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.LEFT
        run3 = p3.add_run()
        set_run_font(run3, date_text, size_pt=14, color_rgb=CATHAY_GREY)

    return cover


# ============================================================================
# 2. CONTENT SLIDE (Layout [4] — left red line + title placeholder)
# ============================================================================

def create_content_slide(prs, title_text=None, topic=None, conclusion=None, layout_idx=4):
    """Create a Layout [4] content slide with dark red title.

    Layout [4] has: left red vertical line (5mm) + title placeholder + content area.
    Title text is set to dark red (#800000) bold.

    Args:
        prs: Presentation
        title_text: simple title string
        topic: topic part (dark red bold)
        conclusion: conclusion part (gold after em-dash)
        layout_idx: layout index (default 4 = "5_Red Slide")
    """
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if topic and conclusion:
        _set_conclusion_title(slide, topic, conclusion)
    elif title_text:
        _set_simple_title(slide, title_text)
    return slide


def create_exec_summary_slide(prs):
    """Create a blank content slide for Executive Summary — no title set yet.
    Title is set manually via _set_simple_title after content is placed."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    return slide


# ============================================================================
# 3. TITLE HELPERS — dark red on title placeholder
# ============================================================================

def _find_title_placeholder(slide):
    """Find the TITLE (type=1) placeholder on a slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.type == 1:
                return shape
    return None


def _set_simple_title(slide, title_text):
    """Set dark red bold title on title placeholder."""
    ph = _find_title_placeholder(slide)
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # Set paragraph spacing
    pPr = p._p.get_or_add_pPr()
    # Remove any existing spacing
    for tag_name in ('a:spcBef', 'a:spcAft'):
        el = pPr.find(qn(tag_name))
        if el is not None:
            pPr.remove(el)

    add_mixed_text(p, title_text, size_pt=TITLE_FONT_PT, bold=True, color_rgb=CATHAY_RED)


def _set_conclusion_title(slide, topic, conclusion):
    """Set conclusion-style title: 'Topic — Conclusion'.
    Topic=dark red bold, separator=gold, Conclusion=gold bold.
    """
    ph = _find_title_placeholder(slide)
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    add_mixed_text(p, topic, size_pt=TITLE_FONT_PT, bold=True, color_rgb=CATHAY_RED)
    add_mixed_text(p, " \u2014 ", size_pt=TITLE_FONT_PT, bold=False, color_rgb=CATHAY_GOLD)
    add_mixed_text(p, conclusion, size_pt=18, bold=True, color_rgb=CATHAY_GOLD)


def set_dark_title(slide, topic, conclusion=None):
    """Public API: set dark red title. With conclusion → Topic — Conclusion in gold."""
    if conclusion:
        _set_conclusion_title(slide, topic, conclusion)
    else:
        _set_simple_title(slide, topic)


# ============================================================================
# 4. SECTION DIVIDER (Layout [11])
# ============================================================================

def create_section_divider(prs, section_title):
    """Create section divider: Layout [11] dark red bg + white left-aligned title."""
    div = prs.slides.add_slide(prs.slide_layouts[11])
    ph = _find_title_placeholder(div)
    if ph is not None:
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        add_mixed_text(p, section_title, size_pt=28, bold=True, color_rgb=CATHAY_WHITE)
    return div


# ============================================================================
# 5. SOURCE FOOTER & PAGE NUMBER
# ============================================================================

def add_source_footer(slide, source_text):
    """Standard source footer — 8pt, y=182mm."""
    txBox = slide.shapes.add_textbox(
        Mm(CL), Mm(SOURCE_Y_MM), Mm(180), Mm(SOURCE_BOX_HEIGHT_MM))
    tf = txBox.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_font(run, f"Source: {source_text}", size_pt=CAPTION_FONT_PT, color_rgb=CATHAY_GREY)
    return txBox


def add_page_number(slide, number, total=None):
    """Add page number at bottom-right, 8pt."""
    text = f"{number}/{total}" if total else str(number)
    txBox = slide.shapes.add_textbox(
        Mm(220), Mm(SOURCE_Y_MM), Mm(24), Mm(SOURCE_BOX_HEIGHT_MM))
    tf = txBox.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    set_run_font(run, text, size_pt=CAPTION_FONT_PT, color_rgb=CATHAY_GREY)
    return txBox


def add_subtitle(slide, text, y_mm=22, size_pt=None):
    """Add a dark grey subtitle line below the title area. (legacy compat)"""
    size_pt = size_pt or SUBTITLE_FONT_PT
    txBox = slide.shapes.add_textbox(Mm(CL), Mm(y_mm), Mm(CW), Mm(8))
    tf = txBox.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_font(run, text, size_pt=size_pt, bold=False, color_rgb=CATHAY_DARK_GREY)
    return txBox


# ============================================================================
# 6. LEGACY ALIASES (backward compat)
# ============================================================================

def set_title(slide, title_text, size_pt=None):
    """Legacy alias for _set_simple_title."""
    _set_simple_title(slide, title_text)


def set_slide_title(slide, title_text, size_pt=None):
    set_title(slide, title_text)


def set_title_with_conclusion(slide, topic, conclusion):
    """Legacy alias for _set_conclusion_title."""
    _set_conclusion_title(slide, topic, conclusion)


def add_dark_title(slide, topic, conclusion=None):
    """Alias for set_dark_title."""
    set_dark_title(slide, topic, conclusion)


__all__ = [
    "create_cover_slide",
    "create_content_slide",
    "create_exec_summary_slide",
    "create_section_divider",
    "set_dark_title",
    "add_dark_title",
    "set_title",
    "set_slide_title",
    "set_title_with_conclusion",
    "add_source_footer",
    "add_page_number",
    "add_subtitle",
]
