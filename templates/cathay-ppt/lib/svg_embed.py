"""
SVG → PNG embedding channel for visual-type slides.

**WHEN TO USE**: 封面、章节 divider、商业模式图、产业链 mapping、复杂 chart。
这些页本质是"成品图"，下游不需要编辑。

**WHEN NOT TO USE**: 财务表格（Comps / Returns / 3-statement）、数字 KPI 卡、正文 bullet 页。
这些页 banker 需要直接双击编辑，走 python-pptx 原生 Table / textbox。

---

## 设计决策

- SVG → **PNG**（不是 EMF，不是原生形状映射）
- PNG 清晰度：**2x-3x DPI**（高清屏投影清晰，不会看到像素）
- 承认"嵌入后不可编辑"的明确代价，保持通道边界清晰

## 后端选择

优先级：
1. `cairosvg`（推荐，rendering 质量最好）
2. `svglib + reportlab`（纯 Python fallback）

安装：
```bash
# 推荐：cairosvg（macOS 需要先装 cairo）
brew install cairo
pip install cairosvg

# 或纯 Python fallback（无需系统依赖）
pip install svglib reportlab
```

## 使用示例

```python
from svg_embed import svg_to_png, embed_svg_slide

# 方式 A：只转 PNG
png_bytes = svg_to_png(svg_string, width_px=1920, height_px=1440)

# 方式 B：直接嵌入 slide（推荐）
from pptx import Presentation
prs = Presentation("template.pptx")
slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

embed_svg_slide(
    slide,
    svg_source=open("cover.svg").read(),  # 或文件路径
    left_in=0, top_in=0,
    width_in=10, height_in=7.5,  # 默认铺满 4:3 slide
)
prs.save("deck.pptx")
```
"""

from __future__ import annotations

import io
import os
from pathlib import Path
from typing import Optional, Union


# ---------------------------------------------------------------------------
# Backend detection
# ---------------------------------------------------------------------------

def _detect_backend() -> str:
    """Return 'cairosvg', 'svglib', or raise."""
    try:
        import cairosvg  # noqa: F401
        return "cairosvg"
    except ImportError:
        pass
    try:
        from svglib.svglib import svg2rlg  # noqa: F401
        from reportlab.graphics import renderPM  # noqa: F401
        return "svglib"
    except ImportError:
        pass
    raise RuntimeError(
        "No SVG renderer available. Install one of:\n"
        "  brew install cairo && pip install cairosvg  (recommended)\n"
        "  pip install svglib reportlab  (pure-python fallback)"
    )


# ---------------------------------------------------------------------------
# Core: SVG string/path -> PNG bytes
# ---------------------------------------------------------------------------

def svg_to_png(
    svg_source: Union[str, bytes, Path],
    width_px: int = 1920,
    height_px: Optional[int] = None,
    dpi: int = 192,
) -> bytes:
    """
    Render an SVG to PNG bytes.

    Parameters
    ----------
    svg_source : str | bytes | Path
        Either raw SVG markup (str/bytes starting with '<') or a file path.
    width_px : int
        Output width in pixels. 1920 at 192 DPI = 10" wide (4:3 slide full width).
    height_px : int | None
        Output height. If None, computed from SVG aspect ratio.
    dpi : int
        Rendering DPI. 192 = 2x retina, good for 4:3 投屏. Use 288 for 3x.

    Returns
    -------
    PNG bytes ready for slide.shapes.add_picture(BytesIO(bytes), ...).
    """
    backend = _detect_backend()

    # Normalize input
    svg_data = _load_svg(svg_source)

    if backend == "cairosvg":
        import cairosvg

        kwargs = {"bytestring": svg_data, "output_width": width_px, "dpi": dpi}
        if height_px is not None:
            kwargs["output_height"] = height_px
        return cairosvg.svg2png(**kwargs)

    elif backend == "svglib":
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM

        # svglib needs a file-like; wrap bytes
        drawing = svg2rlg(io.BytesIO(svg_data))
        if drawing is None:
            raise ValueError("svglib failed to parse SVG")

        # Scale to target width
        scale = width_px / drawing.width
        drawing.width *= scale
        drawing.height *= scale
        drawing.scale(scale, scale)

        buf = io.BytesIO()
        renderPM.drawToFile(drawing, buf, fmt="PNG", dpi=dpi)
        return buf.getvalue()

    raise RuntimeError(f"Unexpected backend: {backend}")


def _load_svg(svg_source: Union[str, bytes, Path]) -> bytes:
    """Normalize str / bytes / path → SVG bytes."""
    if isinstance(svg_source, bytes):
        return svg_source
    if isinstance(svg_source, Path):
        return svg_source.read_bytes()
    if isinstance(svg_source, str):
        stripped = svg_source.lstrip()
        if stripped.startswith("<"):
            return svg_source.encode("utf-8")
        if os.path.exists(svg_source):
            return Path(svg_source).read_bytes()
        raise ValueError(
            "SVG string does not start with '<' and is not a valid file path"
        )
    raise TypeError(f"Unsupported svg_source type: {type(svg_source)}")


# ---------------------------------------------------------------------------
# Convenience: embed into a python-pptx slide
# ---------------------------------------------------------------------------

def embed_svg_slide(
    slide,
    svg_source: Union[str, bytes, Path],
    left_in: float = 0,
    top_in: float = 0,
    width_in: float = 10.0,
    height_in: float = 7.5,
    render_width_px: int = 1920,
    dpi: int = 192,
):
    """
    Render an SVG and insert it as a picture on a python-pptx slide.

    Defaults fill a full 10" x 7.5" 4:3 slide.

    Parameters
    ----------
    slide : pptx.slide.Slide
    svg_source : see svg_to_png
    left_in, top_in, width_in, height_in : float
        Position / size in inches.
    render_width_px : int
        PNG render resolution. Higher = clearer on 大屏投影, larger file.
    dpi : int
        Render DPI. 192 = 2x, 288 = 3x.

    Returns
    -------
    The added Picture shape.
    """
    from pptx.util import Inches

    png_bytes = svg_to_png(svg_source, width_px=render_width_px, dpi=dpi)
    pic = slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        left=Inches(left_in),
        top=Inches(top_in),
        width=Inches(width_in),
        height=Inches(height_in),
    )
    return pic


# ---------------------------------------------------------------------------
# Guard: refuse to embed if the SVG looks like a table
# ---------------------------------------------------------------------------

def _looks_like_table(svg_data: bytes) -> bool:
    """
    Heuristic: if the SVG contains many <rect> in a grid pattern or many <text>
    nodes arranged in rows, warn the caller that this should probably be a
    native PPT Table, not an embedded image.

    Returns True if likely a table (caller should reconsider).
    """
    text_sample = svg_data[:20000].decode("utf-8", errors="ignore").lower()
    # Crude signals
    rect_count = text_sample.count("<rect")
    text_count = text_sample.count("<text")
    return rect_count >= 20 and text_count >= 20


def assert_not_table(svg_source: Union[str, bytes, Path]) -> None:
    """
    Raise if the SVG looks like a data table.

    Call this defensively before embed_svg_slide for data-rich content.
    Caller can suppress by wrapping in try/except, but the default message
    reminds them that banker 需要可编辑 tables.
    """
    svg_data = _load_svg(svg_source)
    if _looks_like_table(svg_data):
        raise ValueError(
            "This SVG looks like a data table (≥20 <rect> + ≥20 <text>).\n"
            "Embedding it as an image makes it uneditable. banker 需要能改数字 — "
            "请改用 python-pptx 原生 Table (lib/text_engine.py::smart_table)。\n"
            "如果确实要嵌入（例如已渲染好的 chart），请 try/except 这条 assert。"
        )


# ---------------------------------------------------------------------------
# Self-test (run module directly)
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    demo_svg = """<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1000 750">
  <rect width="1000" height="750" fill="#FFFFFF"/>
  <rect x="0" y="0" width="6" height="120" fill="#800000"/>
  <text x="40" y="80" font-family="Calibri" font-size="44" font-weight="bold" fill="#800000">
    Project Sealien
  </text>
  <text x="40" y="130" font-family="Calibri" font-size="20" fill="#666666">
    海洋机器人投资推介 · 2026-04
  </text>
</svg>"""

    print(f"Backend: {_detect_backend()}")
    png = svg_to_png(demo_svg, width_px=1920)
    print(f"PNG bytes: {len(png):,}")
    Path("/tmp/svg_embed_demo.png").write_bytes(png)
    print("Wrote /tmp/svg_embed_demo.png")
