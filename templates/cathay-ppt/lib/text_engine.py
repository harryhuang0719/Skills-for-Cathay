"""
Cathay PPT Template — Text Engine & Core Helpers
=================================================
Single-file module containing ALL text fitting functions, layout constants,
font helpers, table/chart helpers, validation, and slide merge logic.

Usage:
    import sys
    sys.path.insert(0, os.path.expanduser('~/.claude/skills/cathay-ppt-template/lib'))
    from text_engine import *
"""

# IRON RULE: ALL lines must be thin rectangles, NEVER connectors.
# Connectors carry corrupt <p:style> references that leak theme XML.
# Use: slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h) with h=Mm(0.2)

import os
import re
import math
import copy
import io
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from constants import *


# ============================================================================
# 1. ANTI-CORRUPTION DEFENSE
# ============================================================================

def _clean_shape(shape):
    """Strip <p:style> XML from shape to prevent theme corruption.
    Call immediately after creating any shape."""
    sp = shape._element
    for pstyle in sp.findall('.//' + qn('p:style')):
        pstyle.getparent().remove(pstyle)


def full_cleanup(pptx_path):
    """Post-save nuclear cleanup: strip ALL <p:style> and theme shadows from PPTX zip.
    Prevents PowerPoint theme corruption caused by python-pptx connector artifacts."""
    import zipfile
    import tempfile
    import shutil

    PSTYLE_RE = re.compile(r'<p:style>.*?</p:style>', re.DOTALL)

    tmp = tempfile.mktemp(suffix='.pptx')
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml') or item.filename.endswith('.rels'):
                    text = data.decode('utf-8')
                    text = PSTYLE_RE.sub('', text)
                    data = text.encode('utf-8')
                zout.writestr(item, data)
    shutil.move(tmp, pptx_path)


def get_char_width(font_pt, is_cjk=False):
    """Get character width in mm for given font size."""
    table = CJK_CHAR_WIDTH if is_cjk else LATIN_CHAR_WIDTH
    pts = sorted(table.keys())
    if font_pt <= pts[0]:
        return table[pts[0]]
    if font_pt >= pts[-1]:
        return table[pts[-1]]
    for i in range(len(pts) - 1):
        if pts[i] <= font_pt <= pts[i + 1]:
            ratio = (font_pt - pts[i]) / (pts[i + 1] - pts[i])
            return table[pts[i]] + ratio * (table[pts[i + 1]] - table[pts[i]])
    return table[10]


# ============================================================================
# 3. CORE TEXT HELPERS
# ============================================================================

def setup_text_frame(tf, word_wrap=True):
    """Apply standard Cathay text frame settings.
    - 0.2cm margins all sides
    - No auto-shrink (MSO_AUTO_SIZE.NONE)
    """
    tf.word_wrap = word_wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left   = MARGIN_ALL
    tf.margin_right  = MARGIN_ALL
    tf.margin_top    = MARGIN_ALL
    tf.margin_bottom = MARGIN_ALL


def format_paragraph(para, indent_left=True, is_bullet=False):
    """Apply standard Cathay paragraph formatting.
    - Spacing before: 4pt, after: 0pt
    - Line spacing: 1.2x (120%)
    - Left indent: 0.5cm (for bulleted paragraphs)
    """
    pPr = para._p.get_or_add_pPr()

    # Spacing before 4pt, after 0pt
    spcBef = pPr.find(qn('a:spcBef'))
    if spcBef is None:
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
    else:
        spcBef.clear()
    spcPts_bef = etree.SubElement(spcBef, qn('a:spcPts'))
    spcPts_bef.set('val', '400')

    spcAft = pPr.find(qn('a:spcAft'))
    if spcAft is None:
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
    else:
        spcAft.clear()
    spcPts_aft = etree.SubElement(spcAft, qn('a:spcPts'))
    spcPts_aft.set('val', '0')

    # Line spacing 1.2x (120%)
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is None:
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
    else:
        lnSpc.clear()
    spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
    spcPct.set('val', str(LINE_SPACING_PCT))

    # Left indent for bulleted paragraphs
    if indent_left and is_bullet:
        pPr.set('indent', str(-Mm(3)))
        pPr.set('marL', str(INDENT_LEFT))


def set_run_font(run, text, size_pt=None, bold=False, color_rgb=None):
    """Set font with auto Chinese/English detection. Default size: 10pt."""
    size_pt = size_pt or DEFAULT_FONT_SIZE
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color_rgb:
        run.font.color.rgb = color_rgb

    has_chinese = any('\u4e00' <= c <= '\u9fff' for c in text)
    if has_chinese:
        run.font.name = "KaiTi"
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:altLang'), 'zh-CN')
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', 'KaiTi')
        # Complex script font (needed for Windows rendering)
        cs = rPr.find(qn('a:cs'))
        if cs is None:
            cs = etree.SubElement(rPr, qn('a:cs'))
        cs.set('typeface', 'KaiTi')
    else:
        run.font.name = "Calibri"
        rPr = run._r.get_or_add_rPr()
        cs = rPr.find(qn('a:cs'))
        if cs is None:
            cs = etree.SubElement(rPr, qn('a:cs'))
        cs.set('typeface', 'Calibri')


def add_mixed_text(para, text, size_pt=None, bold=False, color_rgb=None):
    """Split mixed CJK/Latin text into multiple runs, each with correct font."""
    size_pt = size_pt or DEFAULT_FONT_SIZE
    segments = re.findall(
        r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+|[^\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+',
        text
    )
    for seg in segments:
        if seg.strip() or seg == ' ':
            run = para.add_run()
            set_run_font(run, seg, size_pt=size_pt, bold=bold, color_rgb=color_rgb)


def set_square_bullet(para, color='000000'):
    """Set filled square bullet via PPT buChar XML at 70% size."""
    pPr = para._p.get_or_add_pPr()
    for tag in ('a:buNone', 'a:buChar', 'a:buSzPct', 'a:buClr', 'a:buAutoNum', 'a:buFont'):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    buFont = etree.SubElement(pPr, qn('a:buFont'))
    buFont.set('typeface', 'Calibri')
    buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
    buSzPct.set('val', '70000')
    buClr = etree.SubElement(pPr, qn('a:buClr'))
    srgb = etree.SubElement(buClr, qn('a:srgbClr'))
    srgb.set('val', color)
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', '\u25A0')


def add_bullet_content(tf, items, size_pt=None, color_rgb=None):
    """Add bulleted content to a text frame.
    items: list of (text, level) tuples.
      level 0 = section header (bold red, no bullet, no indent)
      level 1 = bulleted item (filled square, 0.5cm indent)
      level 2+ = sub-item (smaller, grey, 1.0cm indent)
    """
    size_pt = size_pt or DEFAULT_FONT_SIZE
    color = color_rgb or CATHAY_BLACK
    for i, (text, level) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        format_paragraph(p, indent_left=True, is_bullet=(level >= 1))

        run = p.add_run()

        if level == 0:
            set_run_font(run, text, size_pt=size_pt + 2, bold=True, color_rgb=CATHAY_RED)
        elif level == 1:
            set_run_font(run, text, size_pt=size_pt, color_rgb=color)
            set_square_bullet(p)
        elif level >= 2:
            set_run_font(run, text, size_pt=max(size_pt - 1, 8), color_rgb=RGBColor(0x80, 0x80, 0x80))
            set_square_bullet(p, color='808080')
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(Mm(10)))
            pPr.set('indent', str(-Mm(3)))


def add_multi_text(slide, x_mm, y_mm, w_mm, h_mm, segments, align=PP_ALIGN.LEFT, fill_rgb=None):
    """Add a textbox with multiple formatted paragraphs.

    More flexible than add_bullet_content — each segment has its own formatting.

    Args:
        segments: list of (text, kwargs_dict) tuples.
            kwargs: size (pt), bold, color (RGBColor), italic, space_before (pt),
                    space_after (pt), line_spacing (float, e.g. 1.1)
        fill_rgb: optional RGBColor for textbox background

    Returns:
        textbox shape

    Example:
        add_multi_text(slide, CL, CT, CW, 50, [
            ("核心观点", dict(size=14, bold=True, color=CATHAY_RED, space_after=4)),
            ("AI算力需求驱动数据中心建设加速", dict(size=10, color=CATHAY_BLACK)),
            ("预计2030年市场规模达$430B", dict(size=10, color=CATHAY_BLACK, italic=True)),
        ])
    """
    txBox = slide.shapes.add_textbox(Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    tf = txBox.text_frame
    setup_text_frame(tf)

    for idx, (text, kw) in enumerate(segments):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = kw.get('align', align)

        if 'space_before' in kw:
            p.space_before = Pt(kw['space_before'])
        if 'space_after' in kw:
            p.space_after = Pt(kw['space_after'])
        if 'line_spacing' in kw:
            p.line_spacing = kw['line_spacing']

        run = p.add_run()
        set_run_font(
            run, text,
            size_pt=kw.get('size', DEFAULT_FONT_SIZE),
            bold=kw.get('bold', False),
            color_rgb=kw.get('color', CATHAY_BLACK),
        )
        if kw.get('italic', False):
            run.font.italic = True

    if fill_rgb:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill_rgb
    else:
        txBox.fill.background()
    txBox.line.fill.background()

    return txBox


# ============================================================================
# 4. TEXT HEIGHT CALCULATOR
# ============================================================================

def calc_text_height(text_or_paragraphs, box_width_mm, font_pt=10,
                     line_spacing=1.2, margin_mm=4):
    """Calculate rendered text height in mm.

    Args:
        text_or_paragraphs: str or [(text, font_pt, indent_mm), ...]
        box_width_mm: textbox width (mm)
        font_pt: default font size
        line_spacing: line spacing multiplier
        margin_mm: total top+bottom margins (mm)

    Returns:
        float: estimated rendered height (mm)
    """
    usable_w = box_width_mm - margin_mm
    if usable_w <= 0:
        usable_w = 5

    if isinstance(text_or_paragraphs, str):
        paragraphs = [(text_or_paragraphs, font_pt, 0)]
    else:
        paragraphs = text_or_paragraphs

    total_h = margin_mm / 2  # top margin

    for i, (text, p_font, indent_mm) in enumerate(paragraphs):
        if not text.strip():
            total_h += 1.5
            continue

        effective_w = usable_w - indent_mm
        if effective_w <= 0:
            effective_w = usable_w

        segments = re.findall(
            r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+|[^\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+',
            text
        )

        total_text_width = 0
        for seg in segments:
            is_cjk = any('\u4e00' <= c <= '\u9fff' for c in seg)
            char_w = get_char_width(p_font, is_cjk)
            total_text_width += len(seg) * char_w

        n_lines = max(1, math.ceil(total_text_width / effective_w))
        line_h = p_font * 0.3528 * line_spacing

        if i > 0:
            total_h += 1.0

        total_h += n_lines * line_h

    total_h += margin_mm / 2  # bottom margin
    return total_h


def calc_textframe_height(text_frame, box_width_mm):
    """Calculate rendered height from a python-pptx TextFrame object."""
    paragraphs = []
    for p in text_frame.paragraphs:
        p_font = 10
        for r in p.runs:
            if r.font.size:
                p_font = r.font.size / 12700
                break

        indent_mm = 0
        pPr = p._p.get_or_add_pPr()
        marL = pPr.get('marL')
        if marL:
            indent_mm = int(marL) / 36000

        paragraphs.append((p.text, p_font, indent_mm))

    line_sp = 1.2
    if text_frame.paragraphs:
        pPr_el = text_frame.paragraphs[0]._p.find(qn('a:lnSpc'))
        if pPr_el is not None:
            spcPct = pPr_el.find(qn('a:spcPct'))
            if spcPct is not None:
                line_sp = int(spcPct.get('val', '120000')) / 100000

    m_tb = ((text_frame.margin_top or 36000) / 36000 +
            (text_frame.margin_bottom or 36000) / 36000)

    return calc_text_height(paragraphs, box_width_mm, line_spacing=line_sp, margin_mm=m_tb)


# ============================================================================
# 5. SMART CREATORS
# ============================================================================

def smart_textbox(slide, x_mm, y_mm, w_mm, items, max_bottom_mm=180,
                  start_font=10, min_font=8, line_spacing=1.2):
    """Create auto-fit textbox: calculate height, reduce font if overflow.

    Args:
        slide: pptx slide object
        x_mm, y_mm, w_mm: position and width in mm
        items: [(text, level), ...] for add_bullet_content
        max_bottom_mm: bottom limit (mm)
        start_font: initial font size (pt)
        min_font: minimum font size (pt)

    Returns:
        (shape, text_frame, actual_font_pt)
    """
    max_h = max_bottom_mm - y_mm

    chosen_font = start_font
    est_h = max_h + 1  # force at least one iteration
    for try_font in [start_font, start_font - 0.5, start_font - 1,
                     start_font - 1.5, start_font - 2, min_font]:
        if try_font < min_font:
            try_font = min_font

        paras = []
        for text, level in items:
            indent = 5 if level == 1 else (10 if level >= 2 else 0)
            f = try_font + 1 if level == 0 else (try_font - 1 if level >= 2 else try_font)
            paras.append((text, f, indent))

        est_h = calc_text_height(paras, w_mm, line_spacing=line_spacing)

        if est_h <= max_h:
            chosen_font = try_font
            break
        chosen_font = try_font

    actual_h = min(max_h, est_h + 3)
    txBox = slide.shapes.add_textbox(Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(actual_h))
    tf = txBox.text_frame
    setup_text_frame(tf)
    add_bullet_content(tf, items, size_pt=chosen_font)

    return txBox, tf, chosen_font


def smart_table(slide, data, left_mm=None, top_mm=None, width_mm=None,
                max_bottom_mm=180, font_size=9, min_row_h=7):
    """Create auto-fit table: scan cell content, auto-set row heights.

    Returns:
        (table, bottom_y_mm)
    """
    left_mm = left_mm or CL
    top_mm = top_mm or CT
    width_mm = width_mm or CW
    rows = len(data)
    cols = len(data[0]) if data else 0

    col_w = width_mm / cols if cols > 0 else width_mm

    # Calculate row heights based on content
    row_heights = []
    for ri, row_data in enumerate(data):
        max_lines = 1
        for ci, cell_text in enumerate(row_data):
            text = str(cell_text) if cell_text else ""
            if not text:
                continue
            has_cjk = any('\u4e00' <= c <= '\u9fff' for c in text)
            char_w = get_char_width(font_size, has_cjk)
            usable = col_w - 3
            if usable <= 0:
                usable = 5
            lines = max(1, math.ceil(len(text) * char_w / usable))
            lines += text.count('\n')
            max_lines = max(max_lines, lines)

        line_h = font_size * 0.3528 * 1.2
        needed_h = max_lines * line_h + 3
        row_heights.append(max(needed_h, min_row_h))

    total_h = sum(row_heights)

    # Scale down font if table doesn't fit
    if top_mm + total_h > max_bottom_mm:
        for smaller_font in [font_size - 0.5, font_size - 1, font_size - 1.5, font_size - 2]:
            if smaller_font < 7:
                break
            font_size = smaller_font
            row_heights_new = []
            for ri, row_data in enumerate(data):
                max_lines = 1
                for ci, cell_text in enumerate(row_data):
                    text = str(cell_text) if cell_text else ""
                    if not text:
                        continue
                    has_cjk = any('\u4e00' <= c <= '\u9fff' for c in text)
                    char_w = get_char_width(smaller_font, has_cjk)
                    usable = col_w - 3
                    if usable <= 0:
                        usable = 5
                    lines = max(1, math.ceil(len(text) * char_w / usable))
                    max_lines = max(max_lines, lines)
                line_h = smaller_font * 0.3528 * 1.2
                needed_h = max_lines * line_h + 3
                row_heights_new.append(max(needed_h, 6))
            total_h_new = sum(row_heights_new)
            if top_mm + total_h_new <= max_bottom_mm:
                row_heights = row_heights_new
                total_h = total_h_new
                break

    avg_row_h = total_h / rows if rows > 0 else min_row_h
    return add_table(slide, data, left_mm=left_mm, top_mm=top_mm,
                     width_mm=width_mm, row_height=avg_row_h, font_size=font_size)


# ============================================================================
# 6. SLIDE STRUCTURE
# ============================================================================

def create_content_slide(prs, title_text=None, topic=None, conclusion=None):
    """Create a Layout [4] content slide and optionally set its title.

    Args:
        prs: Presentation object
        title_text: simple title (use set_title)
        topic: topic part of conclusion title
        conclusion: conclusion part of conclusion title

    Returns:
        slide object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    if topic and conclusion:
        set_title_with_conclusion(slide, topic, conclusion)
    elif title_text:
        set_title(slide, title_text)
    return slide


def create_cover_slide(prs, fund_name="Cathay Smart Energy Fund",
                       company_name="Company Name", subtitle="Investment Memo",
                       date_text="March 2026"):
    """Create Layout [0] cover slide with standard positioning."""
    cover = prs.slides.add_slide(prs.slide_layouts[0])

    # Fund name
    txBox = cover.shapes.add_textbox(Mm(CL), Mm(51), Mm(200), Mm(15))
    tf = txBox.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_font(run, fund_name, size_pt=20, color_rgb=CATHAY_RED)

    # Company name + subtitle
    txBox2 = cover.shapes.add_textbox(Mm(CL), Mm(76), Mm(216), Mm(30))
    tf2 = txBox2.text_frame
    setup_text_frame(tf2)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    set_run_font(run2, company_name, size_pt=28, bold=True, color_rgb=CATHAY_BLACK)
    p2b = tf2.add_paragraph()
    p2b.alignment = PP_ALIGN.LEFT
    run2b = p2b.add_run()
    set_run_font(run2b, subtitle, size_pt=18, color_rgb=CATHAY_GREY)

    # Date
    txBox3 = cover.shapes.add_textbox(Mm(CL), Mm(127), Mm(76), Mm(10))
    tf3 = txBox3.text_frame
    setup_text_frame(tf3)
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.add_run()
    set_run_font(run3, date_text, size_pt=14, color_rgb=CATHAY_GREY)

    return cover


def set_title(slide, title_text, size_pt=20):
    """Set slide title (white on red bar) -- single-part version."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.type == 1:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                set_run_font(run, title_text, size_pt=size_pt, bold=True,
                             color_rgb=CATHAY_WHITE)
                break


def set_slide_title(slide, title_text, size_pt=20):
    """Alias for set_title."""
    set_title(slide, title_text, size_pt=size_pt)


def set_title_with_conclusion(slide, topic, conclusion):
    """Set slide title: 'Topic -- Conclusion'. Topic=white, Conclusion=gold.

    Use for all Layout [4] content slides.
    Example: set_title_with_conclusion(slide, "投资摘要", "BUY, 目标价$520")
    """
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.type == 1:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run_topic = p.add_run()
                set_run_font(run_topic, topic, size_pt=20, bold=True, color_rgb=CATHAY_WHITE)
                run_sep = p.add_run()
                set_run_font(run_sep, " \u2014 ", size_pt=20, bold=False, color_rgb=CATHAY_GOLD)
                run_conc = p.add_run()
                set_run_font(run_conc, conclusion, size_pt=18, bold=True, color_rgb=CATHAY_GOLD)
                break


def add_subtitle(slide, text, y_mm=22, size_pt=14):
    """Add a gold subtitle line below the title bar."""
    txBox = slide.shapes.add_textbox(Mm(CL), Mm(y_mm), Mm(CW), Mm(8))
    tf = txBox.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_font(run, text, size_pt=size_pt, bold=False, color_rgb=CATHAY_GOLD)
    return txBox


def add_source_footer(slide, source_text):
    """Standard source footer -- 7pt, 5mm height, y=182mm."""
    txBox = slide.shapes.add_textbox(Mm(CL), Mm(SOURCE_Y_MM), Mm(180),
                                     Mm(SOURCE_BOX_HEIGHT_MM))
    tf = txBox.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_font(run, f"Source: {source_text}", size_pt=SOURCE_FONT_PT,
                 color_rgb=CATHAY_GREY)
    return txBox


def add_page_number(slide, number, total=None):
    """Add page number at bottom-right, 7pt."""
    text = f"{number}/{total}" if total else str(number)
    txBox = slide.shapes.add_textbox(Mm(220), Mm(SOURCE_Y_MM), Mm(24),
                                     Mm(SOURCE_BOX_HEIGHT_MM))
    tf = txBox.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    set_run_font(run, text, size_pt=SOURCE_FONT_PT, color_rgb=CATHAY_GREY)
    return txBox


# ============================================================================
# 7. VISUAL ELEMENTS
# ============================================================================

def add_callout_box(slide, x_mm, y_mm, w_mm, h_mm, value, label,
                    bg_color=None, text_color=None):
    """Add a rounded-rect callout box (KPI-style) with value and label."""
    bg_color = bg_color or CATHAY_RED
    text_color = text_color or CATHAY_WHITE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg_color
    sh.line.fill.background()
    tf = sh.text_frame
    setup_text_frame(tf)
    pv = tf.paragraphs[0]
    pv.alignment = PP_ALIGN.CENTER
    rv = pv.add_run()
    set_run_font(rv, value, size_pt=16, bold=True, color_rgb=text_color)
    pl = tf.add_paragraph()
    pl.alignment = PP_ALIGN.CENTER
    rl = pl.add_run()
    set_run_font(rl, label, size_pt=8, color_rgb=text_color)
    return sh


def add_flow_box(slide, x_mm, y_mm, w_mm, h_mm, text, bg_color=None,
                 text_color=None, font_size=9):
    """Add a rounded-rect flow box with centered text."""
    bg_color = bg_color or CATHAY_RED
    text_color = text_color or CATHAY_WHITE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg_color
    sh.line.fill.background()
    tf = sh.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_mixed_text(p, text, size_pt=font_size, bold=True, color_rgb=text_color)
    return sh


def add_arrow(slide, x_mm, y_mm, w_mm=8, h_mm=6, color=None):
    """Add a right-pointing arrow shape."""
    color = color or CATHAY_GOLD
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_down_arrow(slide, x_mm, y_mm, w_mm=6, h_mm=8, color=None):
    """Add a down-pointing arrow shape."""
    color = color or CATHAY_GOLD
    sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_progress_bar(slide, x_mm, y_mm, w_mm, h_mm=4, pct=1.0,
                     fill_color=None, bg_color=None):
    """Add a horizontal progress bar (background + filled portion)."""
    fill_color = fill_color or CATHAY_RED
    bg_color = bg_color or CATHAY_LTGREY

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(bg)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # Filled portion
    fill_w = max(w_mm * pct, 2)
    fg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(fill_w), Mm(h_mm))
    _clean_shape(fg)
    fg.fill.solid()
    fg.fill.fore_color.rgb = fill_color
    fg.line.fill.background()
    return bg, fg


def add_color_block(slide, x_mm, y_mm, w_mm, h_mm, color=None):
    """Add a solid color rectangle block (no border)."""
    color = color or CATHAY_RED
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_kpi_row(slide, kpis, y_mm=None):
    """Add a row of KPI callout boxes.
    kpis: list of (value, label) tuples.
    Returns Y position for content below KPIs.
    """
    y_mm = y_mm or CT
    n = len(kpis)
    bw = (CW - GAP_H * (n - 1)) / n
    for i, (val, lbl) in enumerate(kpis):
        x = X1 + i * (bw + GAP_H)
        add_callout_box(slide, x, y_mm, bw, 22, val, lbl)
    return y_mm + 28


def add_section_marker(slide, x_mm, y_mm, icon_type=None):
    """Place a small colored shape (4x4mm) as a section visual marker."""
    icon_type = icon_type or ICON_INSIGHT
    shape_enum, color_hex = icon_type
    marker = slide.shapes.add_shape(shape_enum, Mm(x_mm), Mm(y_mm), Mm(4), Mm(4))
    _clean_shape(marker)
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    marker.line.fill.background()
    return marker


def auto_assign_icons(items):
    """Auto-assign icons to level-0 items based on keyword matching."""
    icons = {}
    for text, level in items:
        if level != 0:
            continue
        for icon_type, keywords in _ICON_KEYWORD_MAP.items():
            if any(kw in text for kw in keywords):
                icons[text] = icon_type
                break
        else:
            icons[text] = ICON_INSIGHT
    return icons


# ============================================================================
# 8. TABLE HELPERS
# ============================================================================

def add_table(slide, data, left_mm=None, top_mm=None, width_mm=None,
              row_height=None, font_size=None, col_widths=None):
    """Add a Cathay-formatted table with dark-red header and alternating rows.

    Args:
        slide: pptx slide object
        data: 2D list [[row0_col0, row0_col1, ...], [row1_col0, ...], ...]
              First row is treated as header.
        left_mm, top_mm, width_mm: position and size in mm
        row_height: height per row in mm (default 7)
        font_size: font size in pt (default 9)
        col_widths: list of mm widths per column (optional, equal if omitted)

    Returns:
        (table_object, bottom_y_mm)
    """
    left_mm = left_mm or CL
    top_mm = top_mm or CT
    width_mm = width_mm or CW
    row_height = row_height or 7
    font_size = font_size or 9

    rows = len(data)
    cols = len(data[0]) if data else 0
    if rows == 0 or cols == 0:
        return None, top_mm

    total_h = rows * row_height
    table_shape = slide.shapes.add_table(
        rows, cols,
        Mm(left_mm), Mm(top_mm),
        Mm(width_mm), Mm(total_h))
    table = table_shape.table

    # Set column widths
    if col_widths:
        for ci, cw in enumerate(col_widths):
            if ci < cols:
                table.columns[ci].width = Mm(cw)
    else:
        eq_w = width_mm / cols
        for ci in range(cols):
            table.columns[ci].width = Mm(eq_w)

    for i in range(rows):
        table.rows[i].height = Mm(row_height)
        for j in range(cols):
            cell = table.cell(i, j)
            cell.margin_left = Mm(1.5)
            cell.margin_right = Mm(1.5)
            cell.margin_top = Mm(1)
            cell.margin_bottom = Mm(1)

            # Clear default and write text via set_run_font
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            cell_text = str(data[i][j]) if data[i][j] is not None else ""
            run = p.add_run()

            if i == 0:
                set_run_font(run, cell_text, size_pt=font_size, bold=True,
                             color_rgb=CATHAY_WHITE)
            else:
                set_run_font(run, cell_text, size_pt=font_size,
                             color_rgb=CATHAY_BLACK)

            # Background colors
            tcPr = cell._tc.get_or_add_tcPr()
            if i == 0:
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '800000'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)
            elif i % 2 == 0:
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'F2F2F2'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)

    bottom_y = top_mm + total_h
    return table, bottom_y


# ============================================================================
# 9. CHART HELPERS
# ============================================================================

def setup_chart_style():
    """Apply Cathay brand styling to matplotlib (call before creating charts)."""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import matplotlib.font_manager as fm

        available = {f.name for f in fm.fontManager.ttflist}
        font = 'Calibri' if 'Calibri' in available else 'Arial'
        plt.rcParams.update({
            'font.family': font,
            'font.size': 10,
            'axes.prop_cycle': plt.cycler(color=CATHAY_COLORS),
            'axes.edgecolor': '#808080',
            'axes.linewidth': 0.5,
            'grid.color': '#D9D9D9',
            'grid.linewidth': 0.5,
            'figure.facecolor': 'white',
            'axes.facecolor': 'white',
        })
    except ImportError:
        pass


def safe_chart_insert(slide, image_path, x_mm=None, y_mm=None, w_mm=200):
    """Insert chart PNG with width-only sizing, return actual bottom Y (mm).

    Reads actual PNG pixel dimensions, computes rendered height preserving
    aspect ratio. Auto-scales if chart would exceed content zone.

    Returns:
        bottom_y_mm (float): Y coordinate where the chart ends.
    """
    x_mm = x_mm or CL
    y_mm = y_mm or CT

    try:
        from PIL import Image
        with Image.open(image_path) as img:
            px_w, px_h = img.size
    except ImportError:
        # Fallback: assume 16:9 aspect ratio
        px_w, px_h = 1600, 900

    aspect = px_h / px_w
    rendered_h_mm = w_mm * aspect

    _safe_bottom = CB - 3
    bottom_y = y_mm + rendered_h_mm
    if bottom_y > _safe_bottom:
        max_h = _safe_bottom - y_mm
        w_mm = max_h / aspect
        rendered_h_mm = max_h
        bottom_y = y_mm + rendered_h_mm

    slide.shapes.add_picture(image_path, Mm(x_mm), Mm(y_mm), Mm(w_mm))
    return bottom_y


def insert_chart_image(slide, image_path, x_mm=None, y_mm=None, w_mm=None):
    """Insert a chart image (width-only, preserves aspect ratio).
    DEPRECATED: prefer safe_chart_insert().
    """
    x_mm = x_mm or CONTENT_LEFT_CM * 10
    y_mm = y_mm or CONTENT_TOP_CM * 10
    w_mm = w_mm or 200
    slide.shapes.add_picture(image_path, Mm(x_mm), Mm(y_mm), Mm(w_mm))


# ============================================================================
# 10. SAFE LAYOUT HELPERS
# ============================================================================

def safe_textbox(slide, x_mm, y_mm, w_mm, h_mm=None, max_bottom_mm=None):
    """Create a textbox that respects content zone bounds.

    If h_mm is None, fills from y_mm to max_bottom_mm (default 175mm).
    Clamps height to never exceed content zone.

    Returns:
        (shape, text_frame)
    """
    max_bottom_mm = max_bottom_mm or (CB - 3)
    if h_mm is None:
        h_mm = max_bottom_mm - y_mm

    actual_bottom = y_mm + h_mm
    if actual_bottom > max_bottom_mm:
        h_mm = max_bottom_mm - y_mm

    if h_mm <= 0:
        h_mm = 10

    txBox = slide.shapes.add_textbox(Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    tf = txBox.text_frame
    setup_text_frame(tf)
    return txBox, tf


# ============================================================================
# 11. VALIDATION
# ============================================================================

def validate_and_fix(prs):
    """Pre-save validation: check every shape for overflow, auto-reduce font.

    Returns:
        list of fix descriptions
    """
    fixes = []

    for slide in prs.slides:
        for shape in slide.shapes:
            top_mm = shape.top / 36000
            height_mm = shape.height / 36000
            width_mm = shape.width / 36000
            bottom_mm = top_mm + height_mm

            if width_mm < 0.5 or height_mm < 0.5:
                continue

            # Check if it's a footer element
            is_footer = False
            if shape.has_text_frame:
                txt = shape.text_frame.text.lower()
                if 'source:' in txt or (len(txt) < 10 and '/' in txt):
                    is_footer = True

            # Cap shapes that exceed content zone
            if not is_footer and bottom_mm > 181:
                new_h = 181 - top_mm
                if new_h >= 5:
                    shape.height = int(new_h * 36000)
                    fixes.append(f"CAP: {shape.name} bottom {bottom_mm:.0f}->181mm")

            # Check text overflow within textbox
            if shape.has_text_frame and height_mm >= 5:
                est_h = calc_textframe_height(shape.text_frame, width_mm)
                if est_h > height_mm * 1.1:
                    for target in [9, 8.5, 8, 7.5, 7]:
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                if r.font.size and r.font.size / 12700 > target:
                                    r.font.size = Pt(target)
                        new_est = calc_textframe_height(shape.text_frame, width_mm)
                        if new_est <= height_mm:
                            fixes.append(f"FONT: {shape.name} reduced to {target}pt")
                            break

    return fixes


def save_with_validation(prs, path):
    """Save with auto validation + fix + anti-corruption cleanup."""
    fixes = validate_and_fix(prs)
    if fixes:
        print(f"Auto-fixed {len(fixes)} issues before save:")
        for f in fixes[:10]:
            print(f"  {f}")
    prs.save(path)
    full_cleanup(path)
    return fixes


def validate_no_overlap(pptx_path):
    """Check all slides for overlapping shapes. Returns list of issues."""
    prs_check = Presentation(pptx_path)
    issues = []
    for slide_idx, slide in enumerate(prs_check.slides, 1):
        shapes = []
        for sh in slide.shapes:
            l = sh.left / 914400
            t = sh.top / 914400
            r = l + sh.width / 914400
            b = t + sh.height / 914400
            shapes.append((sh.name, l, t, r, b))
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                n1, l1, t1, r1, b1 = shapes[i]
                n2, l2, t2, r2, b2 = shapes[j]
                if l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2:
                    inside = ((l2 >= l1 and r2 <= r1 and t2 >= t1 and b2 <= b1) or
                              (l1 >= l2 and r1 <= r2 and t1 >= t2 and b1 <= b2))
                    if not inside:
                        issues.append(f"Slide {slide_idx}: '{n1}' overlaps '{n2}'")
    return issues


def validate_text_fit(pptx_path):
    """Estimate whether text fits within each textbox. Returns warnings."""
    prs_check = Presentation(pptx_path)
    warnings = []
    for slide_idx, slide in enumerate(prs_check.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            total_text = "".join(p.text for p in tf.paragraphs)
            if not total_text.strip():
                continue
            box_w_mm = shape.width / 36000
            box_h_mm = shape.height / 36000
            has_chinese = any('\u4e00' <= c <= '\u9fff' for c in total_text)
            chars_per_mm = 2.5 if has_chinese else 3.5
            usable_w = box_w_mm - 4
            chars_per_line = max(usable_w * chars_per_mm, 1)
            num_paragraphs = len([p for p in tf.paragraphs if p.text.strip()])
            total_chars = len(total_text)
            est_lines = (total_chars / chars_per_line) + num_paragraphs * 0.3
            line_height_mm = 2.8
            est_height = est_lines * line_height_mm + 4
            if est_height > box_h_mm * 1.15:
                overflow_pct = ((est_height - box_h_mm) / box_h_mm) * 100
                warnings.append(
                    f"Slide {slide_idx}: '{shape.name}' text may overflow "
                    f"(est {est_height:.0f}mm vs box {box_h_mm:.0f}mm, +{overflow_pct:.0f}%)")
    return warnings


def qc_presentation(pptx_path):
    """Run full QC: overlap check + text fit check. Returns (issues, pdf_path)."""
    issues = validate_no_overlap(pptx_path)
    fit_warnings = validate_text_fit(pptx_path)
    if issues:
        print(f"OVERLAP ISSUES ({len(issues)}):")
        for issue in issues:
            print(f"  - {issue}")
    if fit_warnings:
        print(f"TEXT FIT WARNINGS ({len(fit_warnings)}):")
        for w in fit_warnings:
            print(f"  - {w}")
    if not issues and not fit_warnings:
        print("No layout issues found.")
    pdf = export_to_pdf(pptx_path)
    if pdf:
        print(f"PDF exported: {pdf}")
    return issues + fit_warnings, pdf


def export_to_pdf(pptx_path, output_dir=None):
    """Convert PPTX to PDF via LibreOffice for visual QC."""
    output_dir = output_dir or os.path.dirname(pptx_path)
    try:
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf',
             '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=120)
        pdf_path = os.path.join(
            output_dir,
            os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf')
        return pdf_path if os.path.exists(pdf_path) else None
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return None


# ============================================================================
# 12. MERGE
# ============================================================================

def merge_slides(slide_files, output_path, template_path=None, slide_order=None):
    """Merge multiple single-slide files into one deck with image rId remapping.

    Args:
        slide_files: dict {slide_num: path} or list of paths
        output_path: output file path
        template_path: template file path (defaults to TEMPLATE)
        slide_order: list of slide_nums for ordering

    Returns:
        int: number of slides in merged deck
    """
    template_path = template_path or TEMPLATE
    master = Presentation(template_path)

    # Clear template slides
    while len(master.slides) > 0:
        rId = master.slides._sldIdLst[0].rId
        master.part.drop_rel(rId)
        del master.slides._sldIdLst[0]

    if isinstance(slide_files, list):
        slide_files = {i + 1: p for i, p in enumerate(slide_files)}

    if slide_order is None:
        slide_order = sorted(slide_files.keys())

    for src_num in slide_order:
        if src_num not in slide_files:
            continue
        src_path = slide_files[src_num]
        if not os.path.exists(src_path):
            continue

        src_prs = Presentation(src_path)
        src_slide = src_prs.slides[0]

        # Match layout
        layout_name = src_slide.slide_layout.name
        target_layout = master.slide_layouts[4]
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                target_layout = layout
                break

        new_slide = master.slides.add_slide(target_layout)

        # Collect image blobs from source
        img_map = {}
        for rel in src_slide.part.rels.values():
            if "image" in str(rel.reltype):
                try:
                    img_map[rel.rId] = rel.target_part.blob
                except Exception:
                    pass

        # Register images in new slide, build rId mapping
        rId_remap = {}
        for old_rId, blob in img_map.items():
            image_part, new_rId = new_slide.part.get_or_add_image_part(
                io.BytesIO(blob))
            rId_remap[old_rId] = new_rId

        # Copy shapes with remapped image references
        for shape in src_slide.shapes:
            el = copy.deepcopy(shape._element)
            for blip in el.findall('.//' + qn('a:blip')):
                old_rId = blip.get(qn('r:embed'))
                if old_rId in rId_remap:
                    blip.set(qn('r:embed'), rId_remap[old_rId])
            new_slide.shapes._spTree.append(el)

    master.save(output_path)
    return len(master.slides)


def reorder_slides(prs, new_order_1based):
    """Reorder slides in a presentation.

    python-pptx can only append slides; use this to reorder after append.
    new_order_1based: list of 1-based slide indices in desired order.
    Example: [1, 2, 5, 3, 4] moves slide 5 to position 3.
    """
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    if len(new_order_1based) != len(ids):
        raise ValueError(f"new_order has {len(new_order_1based)} items but presentation has {len(ids)} slides")
    reordered = [ids[i - 1] for i in new_order_1based]
    for el in ids:
        sldIdLst.remove(el)
    for el in reordered:
        sldIdLst.append(el)


def clear_slide(slide):
    """Remove all shapes from a slide (for rebuilding content on an existing slide).
    Preserves the slide's layout/master relationship."""
    for shp in list(slide.shapes):
        shp._element.getparent().remove(shp._element)
