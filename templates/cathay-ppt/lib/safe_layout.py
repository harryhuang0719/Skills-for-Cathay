"""
Cathay PPT Template — Safe Layout Helpers
===========================================
safe_textbox() — content-zone-aware textbox creation.

Usage:
    from safe_layout import safe_textbox
"""

from pptx.util import Mm

from constants import (
    CL, CT, CB, CW, CONTENT_BOTTOM_MM,
)
from text_layout import setup_text_frame


def safe_textbox(slide, x_mm, y_mm, w_mm, h_mm=None, max_bottom_mm=None):
    """Create a textbox that respects content zone bounds.

    If h_mm is None, fills from y_mm to max_bottom_mm (default CONTENT_BOTTOM_MM).
    Clamps height to never exceed content zone.

    Returns:
        (shape, text_frame)
    """
    max_bottom_mm = max_bottom_mm or (CB - 3)
    if h_mm is None:
        h_mm = max_bottom_mm - y_mm

    actual_bottom = y_mm + h_mm
    if actual_bottom > max_bottom_mm:
        h_mm = max_bottom_mm - y_mm

    if h_mm <= 0:
        h_mm = 10

    txBox = slide.shapes.add_textbox(Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    tf = txBox.text_frame
    setup_text_frame(tf)
    return txBox, tf


__all__ = [
    "safe_textbox",
]
