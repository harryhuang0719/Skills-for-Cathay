"""
Cathay PPT Template — Pre-built Slide Templates (v2)
======================================================
16 ready-to-use slide template functions using Building Blocks from elements.py.
Each template is 25-40 lines (down from 60-100 in v1).

Usage:
    from slide_templates import template_kpi_dashboard, template_three_column_compare, ...
"""

import os
import tempfile

from pptx.util import Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from constants import (
    CL, CT, CB, CW, CH, GAP_XS, GAP_SM, GAP_MD, GAP_H, GAP_V,
    HALF, THIRD, QUARTER, ONE_THIRD, TWO_THIRDS, ONE_QUARTER, THREE_QUARTER,
    X1, X2_HALF, X2_T23, X2_Q34, X2_MID, X3_RIGHT,
    Y1, Y2_HALF, Y2_MID, Y3_BOT, ROW_HALF, ROW_THIRD, ROW_FULL,
    SOURCE_Y_MM, SOURCE_BOX_HEIGHT_MM, SOURCE_FONT_PT, DEFAULT_FONT_SIZE,
    CONTENT_BOTTOM_MM,
    TITLE_FONT_PT, SUBTITLE_FONT_PT, BODY_FONT_PT, SMALL_FONT_PT, CAPTION_FONT_PT,
    KPI_VALUE_PT, KPI_LABEL_PT,
    CATHAY_RED, CATHAY_GOLD, CATHAY_LTGOLD, CATHAY_ACCENT,
    CATHAY_BLACK, CATHAY_WHITE, CATHAY_GREY, CATHAY_LTGREY, CATHAY_LIGHT_BG,
    CATHAY_VERY_LIGHT, CATHAY_DARK_GREY, CATHAY_DARK_RED,
    ICON_FINANCIAL, ICON_INSIGHT, ICON_RISK, ICON_CATALYST, ICON_ACTION,
)
from fonts import set_run_font, add_mixed_text
from text_layout import setup_text_frame, format_paragraph, add_bullet_content, smart_textbox
from tables import smart_table, add_table
from charts import safe_chart_insert, setup_chart_style
from slides import create_content_slide, create_cover_slide, add_dark_title, set_title, add_source_footer, add_page_number, create_section_divider
from elements import (
    Card, KpiStrip, HeaderBar, ContentPanel, MetricRow,
    SectionBlock, IconBadge, AccentLine, BulletDot, DividerLine,
    add_color_block, add_flow_box, add_arrow, add_down_arrow, add_kpi_row, add_section_marker, auto_assign_icons,
)
from safe_layout import safe_textbox
from merge import _clean_shape


# ============================================================================
# T1: KPI Dashboard — KPI row + bullet content
# ============================================================================

def template_kpi_dashboard(prs, title, subtitle, kpis=None, bullets=None, source=""):
    """KPI row (3-6 boxes) + bullet content below."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    body_y = KpiStrip(slide, X1, CT, kpis or [])
    smart_textbox(slide, X1, body_y, CW, bullets or [],
                  max_bottom_mm=CB, start_font=10, min_font=8)
    add_source_footer(slide, source)
    return slide


# ============================================================================
# T2: Value Chain Flow — horizontal boxes + table
# ============================================================================

def template_value_chain_flow(prs, title, subtitle, chain_items=None, table_data=None, source=""):
    """Horizontal flow chart (N boxes + arrows) + comparison table below."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    chain_items = chain_items or []
    n = len(chain_items) or 1

    flow_h, flow_y, arrow_w = 20, CT + 2, 8
    box_w = (CW - (arrow_w + 3) * (n - 1)) / n

    for i, item_text in enumerate(chain_items):
        cx = X1 + i * (box_w + arrow_w + 3)
        add_flow_box(slide, cx, flow_y, box_w, flow_h, item_text,
                     bg_color=CATHAY_RED, text_color=CATHAY_WHITE, font_size=9)
        if i < n - 1:
            add_arrow(slide, cx + box_w + 1, flow_y + flow_h / 2 - 3,
                      w_mm=arrow_w - 2, h_mm=6, color=CATHAY_GOLD)

    table_top = flow_y + flow_h + GAP_SM + 5
    if table_data:
        smart_table(slide, table_data, top_mm=table_top, max_bottom_mm=CB, font_size=9)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T3: Chart + Analysis — chart on one side, text on the other
# ============================================================================

def template_chart_plus_analysis(prs, title, subtitle, chart_path=None,
                                 analysis_items=None, source="", chart_side="left"):
    """Chart on one side + text analysis on the other."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    chart_w, text_w = HALF, HALF
    chart_x, text_x = (X1, X2_HALF) if chart_side == "left" else (X2_HALF, X1)

    if chart_path and os.path.exists(chart_path):
        safe_chart_insert(slide, chart_path, x_mm=chart_x, y_mm=CT, w_mm=chart_w)

    smart_textbox(slide, text_x, CT, text_w, analysis_items or [],
                  max_bottom_mm=CB, start_font=10, min_font=8)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T4: Comparison Matrix — callouts + table + conclusion
# ============================================================================

def template_comparison_matrix(prs, title, subtitle, callouts=None, table_data=None,
                                conclusion=None, source=""):
    """Callout boxes at top + comparison table + conclusion text."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    cur_y = CT

    if callouts:
        n = len(callouts)
        bw = (CW - GAP_SM * (n - 1)) / n
        for i, (val, lbl) in enumerate(callouts):
            cx = X1 + i * (bw + GAP_SM)
            HeaderBar(slide, cx, cur_y, bw, 18, f"{val}\n{lbl}",
                      color=CATHAY_RED, font_size=10, align=PP_ALIGN.CENTER)
        cur_y += 18 + GAP_SM

    table_bottom = CB if not conclusion else CB - 25
    if table_data:
        _, tbl_bottom = smart_table(slide, table_data, top_mm=cur_y,
                                    max_bottom_mm=table_bottom, font_size=9)
        cur_y = tbl_bottom + GAP_SM

    if conclusion:
        items = [(conclusion, 1)] if isinstance(conclusion, str) else conclusion
        smart_textbox(slide, X1, cur_y, CW, items, max_bottom_mm=CB, start_font=10, min_font=8)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T5: Two-Column Analysis
# ============================================================================

def template_two_column_analysis(prs, title, subtitle, left_items=None, right_items=None,
                                  bottom_kpis=None, source=""):
    """Two-column analysis with optional bottom KPI row."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    content_bottom = CB if not bottom_kpis else CB - 30

    smart_textbox(slide, X1, CT, HALF, left_items or [],
                  max_bottom_mm=content_bottom, start_font=10, min_font=8)
    smart_textbox(slide, X2_HALF, CT, HALF, right_items or [],
                  max_bottom_mm=content_bottom, start_font=10, min_font=8)

    if bottom_kpis:
        KpiStrip(slide, X1, content_bottom + GAP_SM, bottom_kpis)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T6: Sidebar Case Study — 1/4 dark sidebar + 3/4 main
# ============================================================================

def template_sidebar_case_study(prs, title, subtitle, sidebar_metrics=None,
                                 main_items=None, bottom_table=None, source=""):
    """1/4 dark sidebar + 3/4 main content + optional bottom table."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    main_bottom = CB if not bottom_table else CB - 45

    Card(slide, X1, CT, ONE_QUARTER, main_bottom - CT,
         header=None, body=sidebar_metrics,
         color=CATHAY_RED, text_color=CATHAY_WHITE, sidebar_mode=True)

    ContentPanel(slide, X2_Q34, CT, THREE_QUARTER, main_bottom - CT, main_items or [])

    if bottom_table:
        smart_table(slide, bottom_table, top_mm=main_bottom + GAP_SM, max_bottom_mm=CB, font_size=8)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T7: Three-Column Compare
# ============================================================================

def template_three_column_compare(prs, title, subtitle, col1=None, col2=None, col3=None, source=""):
    """Three equal columns with color-accented headers."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    columns = [col1 or {}, col2 or {}, col3 or {}]
    colors = [CATHAY_RED, CATHAY_GOLD, CATHAY_GREY]
    x_positions = [X1, X2_MID, X3_RIGHT]
    header_h = 12

    for i, (col, xp, color) in enumerate(zip(columns, x_positions, colors)):
        header = col.get('header', f'Column {i+1}') if isinstance(col, dict) else f'Column {i+1}'
        HeaderBar(slide, xp, CT, THIRD, header_h, header, color=color, font_size=11,
                  align=PP_ALIGN.CENTER)

        body_y = CT + header_h + GAP_SM
        items = col.get('items', []) if isinstance(col, dict) else col
        if items:
            smart_textbox(slide, xp, body_y, THIRD, items, max_bottom_mm=CB, start_font=10, min_font=8)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T8: Stacked Cases — two cases stacked vertically
# ============================================================================

def template_stacked_cases(prs, title, subtitle, case1=None, case2=None, source=""):
    """Two cases stacked vertically with separator."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    case1 = case1 or {}
    case2 = case2 or {}
    half_h = ROW_HALF
    sep_h = 1

    c1_header = case1.get('header', 'Case 1') if isinstance(case1, dict) else 'Case 1'
    c1_items = case1.get('items', []) if isinstance(case1, dict) else case1
    c1_color = case1.get('color', CATHAY_RED) if isinstance(case1, dict) else CATHAY_RED
    c2_header = case2.get('header', 'Case 2') if isinstance(case2, dict) else 'Case 2'
    c2_items = case2.get('items', []) if isinstance(case2, dict) else case2
    c2_color = case2.get('color', CATHAY_GOLD) if isinstance(case2, dict) else CATHAY_GOLD

    HeaderBar(slide, X1, CT, CW, 10, c1_header, color=c1_color, font_size=11)
    if c1_items:
        smart_textbox(slide, X1, CT + 12, CW, c1_items,
                      max_bottom_mm=CT + half_h - sep_h, start_font=10, min_font=8)

    sep_y = CT + half_h
    add_color_block(slide, X1, sep_y, CW, sep_h, color=CATHAY_LTGREY)

    case2_top = sep_y + sep_h + GAP_SM
    HeaderBar(slide, X1, case2_top, CW, 10, c2_header, color=c2_color, font_size=11)
    if c2_items:
        smart_textbox(slide, X1, case2_top + 12, CW, c2_items,
                      max_bottom_mm=CB, start_font=10, min_font=8)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T9: Risk Cards — severity-colored risk cards
# ============================================================================

def template_risk_cards(prs, title, subtitle, risks=None, source=""):
    """5 risk cards with severity-based colors."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    risks = risks or []
    n = len(risks)

    if n == 0:
        add_source_footer(slide, source)
        return slide

    severity_colors = {
        'high': RGBColor(0xE6, 0x00, 0x00),
        'medium': RGBColor(0xC8, 0xA4, 0x15),
        'low': RGBColor(0x80, 0x80, 0x80),
    }

    if n <= 5:
        card_w = (CW - GAP_SM * (n - 1)) / n
        card_h = CB - CT - GAP_SM

        for idx, risk in enumerate(risks):
            rx = X1 + idx * (card_w + GAP_SM)
            risk_title = risk.get('title', f'Risk {idx+1}') if isinstance(risk, dict) else str(risk)
            risk_desc = risk.get('description', '') if isinstance(risk, dict) else ''
            severity = risk.get('severity') if isinstance(risk, dict) else None
            card_color = severity_colors.get(severity, CATHAY_RED)

            Card(slide, rx, CT, card_w, card_h,
                 header=risk_title, body=[(risk_desc, 1)] if risk_desc else None,
                 color=card_color, header_height=12, start_font=9, min_font=7)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T10: Action Plan — priority funnel + timeline + conclusion
# ============================================================================

def template_action_plan(prs, title, subtitle, tiers=None, timeline=None,
                          conclusion=None, source=""):
    """Priority funnel + timeline + conclusion box."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    tiers = tiers or []
    cur_y = CT
    tier_colors = [CATHAY_RED, CATHAY_GOLD, CATHAY_GREY, CATHAY_LTGREY]

    if tiers:
        n_tiers = len(tiers)
        tier_h = min(18, (CB - CT - 50) / max(n_tiers, 1))

        for i, tier in enumerate(tiers):
            tier_label = tier.get('label', f'Tier {i+1}') if isinstance(tier, dict) else str(tier)
            tier_items = tier.get('items', []) if isinstance(tier, dict) else []
            tier_color = tier.get('color', tier_colors[i % len(tier_colors)]) if isinstance(tier, dict) else tier_colors[i % len(tier_colors)]
            inset = i * 8
            box_w = CW - inset * 2
            box_x = X1 + inset

            text_color = CATHAY_WHITE if i < 2 else CATHAY_BLACK
            items_text = " | ".join(tier_items) if tier_items else ""
            full_text = f"{tier_label}  {items_text}" if items_text else tier_label

            HeaderBar(slide, box_x, cur_y, box_w, tier_h, full_text,
                      color=tier_color, text_color=text_color, font_size=10)
            cur_y += tier_h + 2

    cur_y += GAP_SM

    if timeline:
        n_phases = len(timeline)
        phase_w = (CW - GAP_SM * (n_phases - 1)) / n_phases

        for i, (phase, desc) in enumerate(timeline):
            px = X1 + i * (phase_w + GAP_SM)
            add_flow_box(slide, px, cur_y, phase_w, 10, phase,
                         bg_color=CATHAY_GOLD, text_color=CATHAY_WHITE, font_size=8)
            desc_box, desc_tf = safe_textbox(slide, px, cur_y + 11, phase_w, h_mm=11)
            add_mixed_text(desc_tf.paragraphs[0], desc, size_pt=8, color_rgb=CATHAY_BLACK)

            if i < n_phases - 1:
                add_arrow(slide, px + phase_w + 1, cur_y + 2, w_mm=GAP_SM - 2, h_mm=6, color=CATHAY_GOLD)

        cur_y += 22 + GAP_SM

    if conclusion:
        remaining_h = CB - cur_y
        if remaining_h > 8:
            conclusion_items = [(conclusion, 1)] if isinstance(conclusion, str) else conclusion
            Card(slide, X1, cur_y, CW, remaining_h,
                 header=None, body=conclusion_items,
                 color=CATHAY_GOLD, bg=RGBColor(0xF8, 0xF0, 0xE0))

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T11: Donut Chart — matplotlib donut + insight panel
# ============================================================================

def template_donut_chart(prs, title, subtitle, segments=None, insight_bullets=None, source=""):
    """Matplotlib donut chart + right-side insight panel."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    segments = segments or []
    default_hex = ['800000', 'C8A415', '808080', 'E60000']

    chart_zone_w = CW * 0.55
    chart_img_w = chart_zone_w - 10

    values = [s[1] for s in segments]
    colors_hex = [s[2] if len(s) > 2 and s[2] else default_hex[i % len(default_hex)]
                  for i, s in enumerate(segments)]
    colors_plt = [f'#{c}' for c in colors_hex]

    tmp_path = None
    try:
        setup_chart_style()
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        fig, ax = plt.subplots(figsize=(3.2, 3.2), dpi=150)
        ax.pie(values, colors=colors_plt, startangle=90,
               wedgeprops={'width': 0.35, 'edgecolor': 'white', 'linewidth': 2})
        ax.set_aspect('equal')
        plt.tight_layout(pad=0)

        tmp_path = tempfile.mktemp(suffix='.png')
        fig.savefig(tmp_path, transparent=True, bbox_inches='tight', pad_inches=0.05)
        plt.close(fig)

        donut_x = X1 + (chart_zone_w - chart_img_w) / 2
        chart_bottom = safe_chart_insert(slide, tmp_path, x_mm=donut_x, y_mm=CT + 2, w_mm=chart_img_w)
    except ImportError:
        chart_bottom = CT + 80

    legend_y = chart_bottom + 3
    for i, seg in enumerate(segments):
        label, val = seg[0], seg[1]
        color_hex = seg[2] if len(seg) > 2 and seg[2] else default_hex[i % len(default_hex)]
        r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
        add_color_block(slide, X1 + 5, legend_y, 5, 5, color=RGBColor(r, g, b))
        lbl_box, lbl_tf = safe_textbox(slide, X1 + 12, legend_y, chart_zone_w - 17, h_mm=6)
        set_run_font(lbl_box.text_frame.paragraphs[0].add_run(),
                     f"{label}  {val}%", size_pt=8, color_rgb=CATHAY_BLACK)
        legend_y += 8

    insight_x = X1 + CW * 0.60
    smart_textbox(slide, insight_x, CT, CW * 0.40, insight_bullets or [],
                  max_bottom_mm=CB, start_font=10, min_font=8)

    if tmp_path and os.path.exists(tmp_path):
        try:
            os.remove(tmp_path)
        except OSError:
            pass

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T12: Before / After — side-by-side comparison
# ============================================================================

def template_before_after(prs, title, subtitle, before_items=None, after_items=None, source=""):
    """Side-by-side before/after comparison with visual distinction."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    col_w, header_h = HALF, 14

    HeaderBar(slide, X1, CT, col_w, header_h, "Before / 现状",
              color=CATHAY_LTGREY, text_color=CATHAY_BLACK, align=PP_ALIGN.CENTER)
    HeaderBar(slide, X2_HALF, CT, col_w, header_h, "After / 优化后",
              color=RGBColor(0xFF, 0xF8, 0xE1), text_color=CATHAY_RED, align=PP_ALIGN.CENTER)

    body_y = CT + header_h + GAP_SM
    smart_textbox(slide, X1, body_y, col_w, before_items or [],
                  max_bottom_mm=CB, start_font=10, min_font=8)
    smart_textbox(slide, X2_HALF, body_y, col_w, after_items or [],
                  max_bottom_mm=CB, start_font=10, min_font=8)

    mid_y = CT + (CB - CT) / 2 - 4
    add_arrow(slide, X1 + col_w + 0.5, mid_y, w_mm=GAP_H - 1, h_mm=8, color=CATHAY_RED)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T13: Funnel — progressive narrowing
# ============================================================================

def template_funnel(prs, title, subtitle, stages=None, source=""):
    """Top-to-bottom funnel chart with progressive narrowing."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    stages = stages or []
    n = len(stages)
    if n == 0:
        add_source_footer(slide, source)
        return slide

    funnel_colors = [CATHAY_RED, CATHAY_GOLD, CATHAY_GREY, CATHAY_ACCENT, CATHAY_LTGOLD]
    funnel_zone_w = CW * 0.60
    desc_zone_x = X1 + CW * 0.65
    desc_zone_w = CW * 0.35
    stage_gap = 4
    available_h = CB - CT - GAP_SM
    stage_h = min(22, (available_h - stage_gap * (n - 1)) / n)

    for i, (label, value_str, description) in enumerate(stages):
        shrink_ratio = 1.0 - (i * 0.6 / max(n - 1, 1))
        bar_w = funnel_zone_w * shrink_ratio
        inset = (funnel_zone_w - bar_w) / 2
        bar_x = X1 + inset
        bar_y = CT + i * (stage_h + stage_gap)
        bar_color = funnel_colors[i % len(funnel_colors)]

        add_flow_box(slide, bar_x, bar_y, bar_w, stage_h, "",
                     bg_color=bar_color, text_color=CATHAY_WHITE, font_size=9)

        inner_box, inner_tf = safe_textbox(slide, bar_x + 2, bar_y + 1, bar_w - 4, h_mm=stage_h - 2)
        pi = inner_tf.paragraphs[0]
        pi.alignment = PP_ALIGN.CENTER
        tc = CATHAY_WHITE if i < 3 else CATHAY_BLACK
        set_run_font(pi.add_run(), value_str, size_pt=12, bold=True, color_rgb=tc)
        set_run_font(pi.add_run(), f"  {label}", size_pt=9, color_rgb=tc)

        desc_box, desc_tf = safe_textbox(slide, desc_zone_x, bar_y + 1, desc_zone_w, h_mm=stage_h - 2)
        add_mixed_text(desc_tf.paragraphs[0], description, size_pt=9, color_rgb=CATHAY_BLACK)

        if i < n - 1:
            add_down_arrow(slide, X1 + funnel_zone_w / 2 - 3, bar_y + stage_h,
                           w_mm=6, h_mm=stage_gap, color=CATHAY_GOLD)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T14: SWOT — 2x2 color-coded matrix
# ============================================================================

def template_swot(prs, title, subtitle, strengths=None, weaknesses=None,
                   opportunities=None, threats=None, source=""):
    """Color-coded SWOT 2x2 matrix with Cathay branding."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    quad_gap = 3
    header_h = 16
    quad_w = (CW - quad_gap) / 2
    quad_h = (CB - CT - quad_gap) / 2

    quadrants = [
        ("Strengths",     RGBColor(0xFF, 0xF8, 0xE1), CATHAY_RED,    strengths or [], 0, 0),
        ("Weaknesses",    RGBColor(0xF5, 0xF5, 0xF5), CATHAY_GREY,   weaknesses or [], 1, 0),
        ("Opportunities", RGBColor(0xFF, 0xF8, 0xE1), CATHAY_GOLD,   opportunities or [], 0, 1),
        ("Threats",       RGBColor(0xFF, 0xE0, 0xE0), CATHAY_ACCENT, threats or [], 1, 1),
    ]

    for label, bg_color, hdr_color, items, col, row in quadrants:
        qx = X1 + col * (quad_w + quad_gap)
        qy = CT + row * (quad_h + quad_gap)

        Card(slide, qx, qy, quad_w, quad_h,
             header=label, body=items,
             color=hdr_color, bg=bg_color,
             header_height=header_h, start_font=9, min_font=7)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T15: Waterfall Chart
# ============================================================================

def template_waterfall(prs, title, subtitle, items=None, source=""):
    """Waterfall chart showing incremental changes via matplotlib."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    items = items or []
    tmp_path = None

    try:
        setup_chart_style()
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np

        labels = [it[0] for it in items]
        values__ = [it[1] for it in items]
        is_totals = [it[2] if len(it) > 2 else False for it in items]

        n = len(items)
        cumulative = 0
        bottoms, bar_vals, bar_colors = [], [], []

        for i in range(n):
            val = values__[i]
            is_total = is_totals[i]
            if is_total:
                bottoms.append(0)
                bar_vals.append(cumulative)
                bar_colors.append('#800000')
            else:
                if val >= 0:
                    bottoms.append(cumulative)
                    bar_vals.append(val)
                    bar_colors.append('#C8A415')
                    cumulative += val
                else:
                    cumulative += val
                    bottoms.append(cumulative)
                    bar_vals.append(abs(val))
                    bar_colors.append('#E60000')

        fig, ax = plt.subplots(figsize=(7, 3.5), dpi=150)
        x_pos = np.arange(n)
        bars = ax.bar(x_pos, bar_vals, bottom=bottoms, color=bar_colors,
                      edgecolor='white', linewidth=0.8, width=0.6)

        max_val = max(abs(v) for v in values__) if values__ else 1
        for i, (bar, bv, bt) in enumerate(zip(bars, bar_vals, bottoms)):
            val = values__[i]
            label_y = bt + bv + (max_val * 0.02)
            display_text = f'{val:g}' if is_totals[i] else f'{val:+g}'
            ax.text(bar.get_x() + bar.get_width() / 2, label_y,
                    display_text, ha='center', va='bottom', fontsize=8, fontweight='bold')

        for i in range(n - 1):
            if not is_totals[i]:
                top_of_current = bottoms[i] + bar_vals[i]
                ax.plot([x_pos[i] + 0.3, x_pos[i + 1] - 0.3],
                        [top_of_current, top_of_current],
                        color='#808080', linewidth=0.5, linestyle='--')

        ax.set_xticks(x_pos)
        ax.set_xticklabels(labels, fontsize=8, rotation=30, ha='right')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_linewidth(0.5)
        ax.spines['bottom'].set_linewidth(0.5)
        ax.tick_params(axis='y', labelsize=8)
        ax.grid(axis='y', alpha=0.3)
        plt.tight_layout()

        tmp_path = tempfile.mktemp(suffix='.png')
        fig.savefig(tmp_path, transparent=False, bbox_inches='tight', pad_inches=0.1, facecolor='white')
        plt.close(fig)

        safe_chart_insert(slide, tmp_path, x_mm=X1, y_mm=CT, w_mm=CW)
    except ImportError:
        fallback_items = []
        for i in range(len(items)):
            lbl, val = items[i][0], items[i][1]
            is_total = items[i][2] if len(items[i]) > 2 else False
            prefix = "TOTAL: " if is_total else ("+" if val >= 0 else "")
            fallback_items.append((f"{lbl}: {prefix}{val}", 1))
        smart_textbox(slide, X1, CT, CW, fallback_items, max_bottom_mm=CB, start_font=10, min_font=8)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try:
                os.remove(tmp_path)
            except OSError:
                pass

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T16: Stakeholder Map
# ============================================================================

def template_stakeholder_map(prs, title, subtitle, stakeholders=None, source=""):
    """Visual stakeholder / relationship map with center + surrounding nodes."""
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    stakeholders = stakeholders or []
    map_cx = CL + CW / 2
    map_cy = CT + CH / 2

    pos_map = {
        'center':       (map_cx, map_cy),
        'top':          (map_cx, CT + 20),
        'bottom':       (map_cx, CB - 20),
        'left':         (CL + 30, map_cy),
        'right':        (CL + CW - 30, map_cy),
        'top-left':     (CL + 40, CT + 25),
        'top-right':    (CL + CW - 40, CT + 25),
        'bottom-left':  (CL + 40, CB - 25),
        'bottom-right': (CL + CW - 40, CB - 25),
    }

    center_sh = None
    outer_shs = []
    for name, role, position in stakeholders:
        if position == 'center':
            center_sh = (name, role, position)
        else:
            outer_shs.append((name, role, position))

    if center_sh is None and stakeholders:
        center_sh = (stakeholders[0][0], stakeholders[0][1], 'center')
        outer_shs = list(stakeholders[1:])

    center_d = 25
    outer_d = 18

    # Connecting lines
    if center_sh:
        c_x, c_y = pos_map['center']
        for name, role, position in outer_shs:
            pos_key = position if position in pos_map else 'right'
            o_x, o_y = pos_map[pos_key]
            if abs(o_x - c_x) >= abs(o_y - c_y):
                mid_y = (c_y + o_y) / 2
                line_x = min(c_x, o_x)
                add_color_block(slide, line_x, mid_y - 0.5, abs(o_x - c_x), 1, color=CATHAY_LTGREY)
            else:
                mid_x = (c_x + o_x) / 2
                line_y = min(c_y, o_y)
                add_color_block(slide, mid_x - 0.5, line_y, 1, abs(o_y - c_y), color=CATHAY_LTGREY)

    # Center node
    if center_sh:
        c_name, c_role, _ = center_sh
        c_x, c_y = pos_map['center']
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Mm(c_x - center_d / 2), Mm(c_y - center_d / 2),
                                    Mm(center_d), Mm(center_d))
        _clean_shape(sh)
        sh.fill.solid()
        sh.fill.fore_color.rgb = CATHAY_RED
        sh.line.fill.background()

        tf = sh.text_frame
        setup_text_frame(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run_font(p.add_run(), c_name, size_pt=9, bold=True, color_rgb=CATHAY_WHITE)
        if c_role:
            pr = tf.add_paragraph()
            pr.alignment = PP_ALIGN.CENTER
            set_run_font(pr.add_run(), c_role, size_pt=7, color_rgb=CATHAY_WHITE)

    # Outer nodes
    for name, role, position in outer_shs:
        pos_key = position if position in pos_map else 'right'
        o_x, o_y = pos_map[pos_key]
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Mm(o_x - outer_d / 2), Mm(o_y - outer_d / 2),
                                    Mm(outer_d), Mm(outer_d))
        _clean_shape(sh)
        sh.fill.solid()
        sh.fill.fore_color.rgb = CATHAY_GOLD
        sh.line.fill.background()

        tf = sh.text_frame
        setup_text_frame(tf)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run_font(p.add_run(), name, size_pt=8, bold=True, color_rgb=CATHAY_WHITE)
        if role:
            pr = tf.add_paragraph()
            pr.alignment = PP_ALIGN.CENTER
            set_run_font(pr.add_run(), role, size_pt=6, color_rgb=CATHAY_WHITE)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T17: Timeline — milestone timeline (NEW in v2)
# ============================================================================

def template_timeline(prs, title, subtitle, milestones=None, source=""):
    """Horizontal milestone timeline with event cards.

    Args:
        milestones: list of dicts:
            { 'date': "2023", 'title': "Founded", 'desc': "Company incorporated..." }
    """
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    milestones = milestones or []
    n = len(milestones)

    if n == 0:
        add_source_footer(slide, source)
        return slide

    # Horizontal timeline line at mid-Y
    line_y = CT + CH / 2
    add_color_block(slide, X1, line_y, CW, 1, color=CATHAY_LTGREY)

    card_w = min((CW - GAP_MD * (n - 1)) / n, 80)
    card_h = 40

    for i, m in enumerate(milestones):
        cx = X1 + i * (card_w + GAP_MD)
        # Alternate above/below the timeline
        if i % 2 == 0:
            cy = line_y - card_h - GAP_SM
        else:
            cy = line_y + GAP_SM

        date = m.get('date', '') if isinstance(m, dict) else str(m)
        m_title = m.get('title', '') if isinstance(m, dict) else ''
        desc = m.get('desc', '') if isinstance(m, dict) else ''

        # Dot on timeline
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Mm(cx + card_w / 2 - 3), Mm(line_y - 3),
                                     Mm(6), Mm(6))
        _clean_shape(dot)
        dot.fill.solid()
        dot.fill.fore_color.rgb = CATHAY_RED
        dot.line.fill.background()

        # Card
        Card(slide, cx, cy, card_w, card_h,
             header=date, body=[(m_title, 0), (desc, 1)] if m_title else [(desc, 1)],
             color=CATHAY_RED, header_height=10, start_font=8, min_font=7)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T19: Number Story — big KPI numbers with short explanations (NEW in v2)
# ============================================================================

def template_number_story(prs, title, subtitle, metrics=None, source=""):
    """3-4 oversized KPI cards with short narrative explanations.

    Args:
        metrics: list of dicts:
            { 'value': "$430B", 'label': "TAM 2028E", 'insight': "Global data center market..." }
    """
    slide = create_content_slide(prs, topic=title, conclusion=subtitle)
    metrics = metrics or []
    n = len(metrics)

    if n == 0:
        add_source_footer(slide, source)
        return slide

    card_w = (CW - GAP_MD * (n - 1)) / n
    card_h = CB - CT - GAP_SM

    for i, m in enumerate(metrics):
        cx = X1 + i * (card_w + GAP_MD)
        value = m.get('value', '') if isinstance(m, dict) else str(m)
        label = m.get('label', '') if isinstance(m, dict) else ''
        insight = m.get('insight', '') if isinstance(m, dict) else ''

        body_items = [(value, 0), (label, 1)]
        if insight:
            body_items.append((insight, 2))

        Card(slide, cx, CT, card_w, card_h,
             header=None, body=body_items,
             color=CATHAY_RED, start_font=10, min_font=8)

    add_source_footer(slide, source)
    return slide


# ============================================================================
# T20: Executive Summary — lead-in + highlights + transaction summary (NEW v3)
# ============================================================================

def template_executive_summary(prs, title="Executive Summary", lead_in=None,
                                highlights=None, transaction=None, source=""):
    """Executive summary page: lead-in paragraph + key highlights + transaction summary.

    Args:
        prs: Presentation object
        title: page title (default "Executive Summary")
        lead_in: opening paragraph text (str)
        highlights: list of highlight strings
        transaction: dict with keys 'structure', 'valuation', 'amount', 'use_of_funds'
        source: source text for footer

    Returns:
        slide object
    """
    slide = create_content_slide(prs, title_text=title)

    # Upper region: Lead-in paragraph
    lead_y = CT
    lead_items = [(lead_in, 1)] if lead_in else []
    lead_bottom = CT + 25  # ~25mm for lead-in

    if lead_items:
        smart_textbox(slide, X1, lead_y, CW, lead_items,
                      max_bottom_mm=lead_bottom, start_font=BODY_FONT_PT + 0.5, min_font=BODY_FONT_PT)

    # Middle region: Key Highlights (left 55%) + Transaction Summary (right 45%)
    hl_y = lead_bottom + GAP_SM

    # Divider before highlights
    DividerLine(slide, X1, hl_y, CW, color=CATHAY_LTGREY)

    hl_header_y = hl_y + GAP_SM
    SectionBlock(slide, X1, hl_header_y, CW * 0.55 - GAP_XS, "Key Highlights",
                 color=CATHAY_RED, font_size=SUBTITLE_FONT_PT - 2, height=12)

    if highlights:
        hl_items = [(h, 1) for h in highlights]
        smart_textbox(slide, X1, hl_header_y + 14, CW * 0.55 - GAP_XS, hl_items,
                      max_bottom_mm=CB - 5, start_font=BODY_FONT_PT, min_font=SMALL_FONT_PT)

    # Right: Transaction Summary box
    tx_x = X1 + CW * 0.58
    tx_w = CW * 0.42
    if transaction:
        SectionBlock(slide, tx_x, hl_header_y, tx_w, "Transaction Summary",
                     color=CATHAY_DARK_RED, font_size=SUBTITLE_FONT_PT - 2, height=12)

        tx_items = []
        if isinstance(transaction, dict):
            for key, label in [('structure', 'Structure'), ('valuation', 'Valuation'),
                               ('amount', 'Amount'), ('use_of_funds', 'Use of Funds')]:
                if key in transaction:
                    tx_items.append((f"{label}: {transaction[key]}", 1))
        else:
            tx_items = [(str(transaction), 1)]

        if tx_items:
            Card(slide, tx_x, hl_header_y + 14, tx_w, 60,
                 header=None, body=tx_items,
                 color=CATHAY_GOLD, bg=RGBColor(0xFF, 0xF8, 0xE1),
                 start_font=SMALL_FONT_PT, min_font=SMALL_FONT_PT - 1)

    add_source_footer(slide, source)
    return slide


__all__ = [
    # Original 16 templates
    "template_kpi_dashboard",
    "template_value_chain_flow",
    "template_chart_plus_analysis",
    "template_comparison_matrix",
    "template_two_column_analysis",
    "template_sidebar_case_study",
    "template_three_column_compare",
    "template_stacked_cases",
    "template_risk_cards",
    "template_action_plan",
    "template_donut_chart",
    "template_before_after",
    "template_funnel",
    "template_swot",
    "template_waterfall",
    "template_stakeholder_map",
    # New in v2
    "template_timeline",
    "template_number_story",
    # New in v3
    "template_executive_summary",
]
