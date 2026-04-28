"""
Cathay PPT Template — Building Blocks
=======================================
Composable visual primitives that eliminate repeated boilerplate in slide templates.
Each block returns its bottom_y_mm for chainable layout.

Blocks:
    HeaderBar     — colored title bar with white text
    ContentPanel  — text panel with optional background
    KpiStrip      — horizontal row of KPI callout boxes
    Card          — complete card: HeaderBar + bullet body + optional background
    MetricRow     — row of small metric cards with value+label+trend

Usage:
    from elements import Card, KpiStrip, HeaderBar, ContentPanel, MetricRow
"""

from pptx.util import Mm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from constants import (
    CW, GAP_XS, GAP_SM, GAP_MD, GAP_LG, GAP_H, GAP_V,
    CL, CT, CB, CH, X1,
    CATHAY_RED, CATHAY_GOLD, CATHAY_WHITE, CATHAY_BLACK, CATHAY_GREY,
    CATHAY_LTGREY, CATHAY_LIGHT_BG, CATHAY_VERY_LIGHT, CATHAY_SOFT_PINK, CATHAY_DARK_GREY,
    DEFAULT_FONT_SIZE, BODY_FONT_PT, SMALL_FONT_PT, CAPTION_FONT_PT, KPI_VALUE_PT, KPI_LABEL_PT,
    TITLE_FONT_PT, SUBTITLE_FONT_PT,
)
from fonts import set_run_font, add_mixed_text
from text_layout import setup_text_frame, format_paragraph, add_bullet_content, smart_textbox
from merge import _clean_shape
from safe_layout import safe_textbox


# ============================================================================
# 1. HeaderBar — colored title bar with white text
# ============================================================================

def HeaderBar(slide, x, y, w, h, title, color=None, text_color=None, font_size=11,
              align=PP_ALIGN.LEFT):
    """Colored title bar: solid rectangle + white text overlay.

    Args:
        slide: pptx slide object
        x, y, w, h: position and size in mm
        title: text string
        color: bar background RGBColor (default CATHAY_RED)
        text_color: text RGBColor (default CATHAY_WHITE)
        font_size: font size in pt (default 11)
        align: text alignment (default LEFT)

    Returns:
        bottom_y_mm (float)
    """
    color = color or CATHAY_RED
    text_color = text_color or CATHAY_WHITE

    # Background rectangle
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(x), Mm(y), Mm(w), Mm(h))
    _clean_shape(bar)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Text overlay
    if h >= 6:
        # Enough height for text
        hdr_box, hdr_tf = safe_textbox(slide, x + GAP_XS, y + 1, w - GAP_XS * 2, h_mm=h - 2)
        p = hdr_tf.paragraphs[0]
        p.alignment = align
        add_mixed_text(p, title, size_pt=font_size, bold=True, color_rgb=text_color)
    else:
        # Very thin bar — text outside
        hdr_box, hdr_tf = safe_textbox(slide, x + GAP_XS, y + 0.5, w - GAP_XS * 2, h_mm=5)
        p = hdr_tf.paragraphs[0]
        p.alignment = align
        add_mixed_text(p, title, size_pt=min(font_size, 8), bold=True, color_rgb=color)

    return y + h


# ============================================================================
# 2. ContentPanel — text panel with optional background
# ============================================================================

def ContentPanel(slide, x, y, w, h, items, bg=None, start_font=10, min_font=8):
    """Text content panel: optional background fill + smart_textbox.

    Args:
        slide: pptx slide object
        x, y, w, h: position and size in mm
        items: list of (text, level) tuples for add_bullet_content
        bg: optional RGBColor for panel background (default transparent)
        start_font: initial font size (default 10)
        min_font: minimum font size (default 8)

    Returns:
        bottom_y_mm (float)
    """
    if bg:
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(x), Mm(y), Mm(w), Mm(h))
        _clean_shape(panel)
        panel.fill.solid()
        panel.fill.fore_color.rgb = bg
        panel.line.fill.background()

    max_bottom = y + h
    _, _, _ = smart_textbox(
        slide, x + GAP_XS, y + GAP_XS, w - GAP_XS * 2,
        items, max_bottom_mm=max_bottom - GAP_XS,
        start_font=start_font, min_font=min_font
    )

    return y + h


# ============================================================================
# 3. KpiStrip — horizontal row of KPI callout boxes
# ============================================================================

def KpiStrip(slide, x, y, kpis, width=None, box_h=22, text_color=None):
    """Horizontal row of KPI callout boxes.

    Args:
        slide: pptx slide object
        x, y: starting position in mm
        kpis: list of (value, label) tuples, e.g. [("$1.2B", "Revenue"), ("45%", "CAGR")]
        width: total strip width (default CW, i.e. full content width)
        box_h: height of each KPI box in mm (default 22)
        text_color: text color (default CATHAY_WHITE)

    Returns:
        bottom_y_mm (float)
    """
    width = width or CW
    text_color = text_color or CATHAY_WHITE
    n = len(kpis)
    bw = (width - GAP_SM * (n - 1)) / n

    for i, (val, lbl) in enumerate(kpis):
        bx = x + i * (bw + GAP_SM)
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Mm(bx), Mm(y), Mm(bw), Mm(box_h))
        _clean_shape(sh)
        sh.fill.solid()
        sh.fill.fore_color.rgb = CATHAY_RED
        sh.line.fill.background()
    tf = sh.text_frame
    setup_text_frame(tf)
    pv = tf.paragraphs[0]
    pv.alignment = PP_ALIGN.CENTER
    rv = pv.add_run()
    set_run_font(rv, val, size_pt=KPI_VALUE_PT, bold=True, color_rgb=text_color)
    pl = tf.add_paragraph()
    pl.alignment = PP_ALIGN.CENTER
    rl = pl.add_run()
    set_run_font(rl, lbl, size_pt=KPI_LABEL_PT, color_rgb=text_color)

    return y + box_h + GAP_SM


# ============================================================================
# 4. Card — complete card component (header bar + body + accent border)
# ============================================================================

def Card(slide, x, y, w, h, header=None, body=None, color=None,
         text_color=None, bg=None, sidebar_mode=False, header_height=12,
         start_font=10, min_font=8):
    """Complete card component: accent left border + header bar + body content.

    Args:
        slide: pptx slide object
        x, y, w, h: card position and size in mm
        header: header title string (optional — no header if None)
        body: list of (text, level) tuples for bullet content
        color: accent color (header bar + left border). Default CATHAY_RED.
        text_color: header text color. Default CATHAY_WHITE.
        bg: card body background color. Default CATHAY_VERY_LIGHT.
        sidebar_mode: if True, the entire card is filled with the accent color
                      (use for 1/4 dark sidebar pattern)
        header_height: height of header bar in mm (default 12)
        start_font: initial body font size (default 10)
        min_font: minimum body font size (default 8)

    Returns:
        bottom_y_mm (float)
    """
    color = color or CATHAY_RED
    text_color = text_color or CATHAY_WHITE
    bg = bg or CATHAY_VERY_LIGHT

    if sidebar_mode:
        # Full-color sidebar: solid fill + white text directly
        sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(x), Mm(y), Mm(w), Mm(h))
        _clean_shape(sidebar)
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = color
        sidebar.line.fill.background()

        if body:
            ContentPanel(slide, x + GAP_XS, y + GAP_XS, w - GAP_XS * 2, h - GAP_XS * 2,
                         body, bg=None, start_font=start_font, min_font=min_font)
        return y + h

    # Normal card mode
    rem_h = h
    cur_y = y

    # Thin left accent border (2mm)
    border_w = 2
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Mm(x), Mm(y), Mm(border_w), Mm(h))
    _clean_shape(border)
    border.fill.solid()
    border.fill.fore_color.rgb = color
    border.line.fill.background()

    # Content area starts after the left border
    content_x = x + border_w + GAP_XS
    content_w = w - border_w - GAP_XS

    if header:
        HeaderBar(slide, content_x, cur_y, content_w, header_height, header,
                  color=color, text_color=text_color, font_size=10)
        cur_y += header_height + GAP_XS
        rem_h = h - (cur_y - y)

    # Body background fill
    if rem_h > 4 and bg:
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Mm(content_x), Mm(cur_y),
                                          Mm(content_w), Mm(rem_h))
        _clean_shape(bg_shape)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = bg
        bg_shape.line.fill.background()

    # Body content
    if body and rem_h > 8:
        ContentPanel(slide, content_x, cur_y, content_w, rem_h, body,
                     bg=None, start_font=start_font, min_font=min_font)

    return y + h


# ============================================================================
# 5. MetricRow — row of small metric cards
# ============================================================================

def MetricRow(slide, x, y, metrics, width=None, card_h=32, inner_gap=GAP_SM):
    """Horizontal row of small metric cards. Each card shows value + label + optional trend arrow.

    Args:
        slide: pptx slide object
        x, y: starting position in mm
        metrics: list of dicts:
            { 'value': "$1.2B", 'label': "Revenue", 'trend': "up", 'sub': "+15% YoY" }
            trend: None / "up" / "down" / "flat"
            sub: optional subtitle line
        width: total row width (default CW)
        card_h: height per metric card (default 32)
        inner_gap: gap between cards (default GAP_SM)

    Returns:
        bottom_y_mm (float)
    """
    width = width or CW
    n = len(metrics)
    card_w = (width - inner_gap * (n - 1)) / n

    for i, m in enumerate(metrics):
        cx = x + i * (card_w + inner_gap)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Mm(cx), Mm(y), Mm(card_w), Mm(card_h))
        _clean_shape(card)
        card.fill.solid()
        card.fill.fore_color.rgb = CATHAY_VERY_LIGHT
        card.line.color.rgb = CATHAY_LTGREY
        card.line.width = Mm(0.3)

        # Text overlay
        inner_box, inner_tf = safe_textbox(slide, cx + 2, y + 2,
                                           card_w - 4, h_mm=card_h - 4)

        # Value line
        pv = inner_tf.paragraphs[0]
        pv.alignment = PP_ALIGN.CENTER
        val_text = m.get('value', '')
        trend = m.get('trend')
        if trend == 'up':
            val_text = '\u25b2 ' + val_text
        elif trend == 'down':
            val_text = '\u25bc ' + val_text
        set_run_font(pv.add_run(), val_text, size_pt=KPI_VALUE_PT - 4, bold=True, color_rgb=CATHAY_RED)

        # Label line
        if m.get('label'):
            pl = inner_tf.add_paragraph()
            pl.alignment = PP_ALIGN.CENTER
            set_run_font(pl.add_run(), m['label'], size_pt=KPI_LABEL_PT, color_rgb=CATHAY_GREY)

        # Subtitle line
        if m.get('sub'):
            ps = inner_tf.add_paragraph()
            ps.alignment = PP_ALIGN.CENTER
            set_run_font(ps.add_run(), m['sub'], size_pt=CAPTION_FONT_PT, color_rgb=CATHAY_DARK_GREY)

    return y + card_h + GAP_SM


# ============================================================================
# 6. SectionBlock — full-width colored section divider within a slide
# ============================================================================

def SectionBlock(slide, x, y, w, title, color=None, text_color=None, font_size=12, height=14):
    """Full-width colored block with white text — use as in-slide section header.

    Distinct from HeaderBar (which is smaller and used for column/card headers).
    SectionBlock is for major content partitions within a single slide.

    Args:
        slide: pptx slide object
        x, y, w, height: position and size in mm
        title: section heading text
        color: block background (default CATHAY_RED)
        text_color: text color (default CATHAY_WHITE)
        font_size: title font size (default 12)
        height: block height (default 14)

    Returns:
        bottom_y_mm (float)
    """
    color = color or CATHAY_RED
    text_color = text_color or CATHAY_WHITE

    # Background rectangle
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(x), Mm(y), Mm(w), Mm(height))
    _clean_shape(block)
    block.fill.solid()
    block.fill.fore_color.rgb = color
    block.line.fill.background()

    # White text overlay
    if height >= 8:
        lbl_box, lbl_tf = safe_textbox(slide, x + GAP_SM, y + 2, w - GAP_SM * 2, h_mm=height - 4)
        p = lbl_tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        add_mixed_text(p, title, size_pt=font_size, bold=True, color_rgb=text_color)

    return y + height


# ============================================================================
# 7. IconBadge — small colored badge with icon + label
# ============================================================================

def IconBadge(slide, x, y, label, icon_shape=None, color=None, text_color=None, size=8, font_size=7):
    """Small colored badge (rounded rect) with icon marker and label text.

    Args:
        slide: pptx slide object
        x, y: position in mm
        label: badge text
        icon_shape: MSO_SHAPE for icon marker (default OVAL dot)
        color: badge background (default CATHAY_RED)
        text_color: text color (default CATHAY_WHITE)
        size: badge height in mm (default 8)
        font_size: label font size (default 7)

    Returns:
        bottom_y_mm (float)
    """
    from pptx.util import Pt as _Pt

    color = color or CATHAY_RED
    text_color = text_color or CATHAY_WHITE

    # Estimate text width for badge sizing
    char_w = 2.0 if all(ord(c) < 128 for c in label) else 3.2
    badge_w = max(len(label) * char_w + 10, 16)

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Mm(x), Mm(y), Mm(badge_w), Mm(size))
    _clean_shape(badge)
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()

    # Label text
    lbl_box, lbl_tf = safe_textbox(slide, x + 2, y + 1, badge_w - 4, h_mm=size - 2)
    p = lbl_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_mixed_text(p, label, size_pt=font_size, bold=True, color_rgb=text_color)

    return y + size


# ============================================================================
# 8. AccentLine — thin colored vertical accent bar
# ============================================================================

def AccentLine(slide, x, y, h, color=None, width=2):
    """Thin vertical accent bar for visual emphasis beside text.

    Use at the left edge of callout boxes or key insight panels.

    Args:
        slide: pptx slide object
        x, y: position in mm
        h: height in mm
        color: bar color (default CATHAY_RED)
        width: bar width in mm (default 2)

    Returns:
        bottom_y_mm (float)
    """
    color = color or CATHAY_RED
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Mm(x), Mm(y), Mm(width), Mm(h))
    _clean_shape(bar)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return y + h


# ============================================================================
# 9. BulletDot — small colored circle marker
# ============================================================================

def BulletDot(slide, x, y, color=None, diameter=4):
    """Small colored circle for visual bullet markers or status indicators.

    Args:
        slide: pptx slide object
        x, y: center position in mm
        color: dot color (default CATHAY_RED)
        diameter: dot diameter in mm (default 4)

    Returns:
        diameter value (for spacing)
    """
    color = color or CATHAY_RED
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Mm(x - diameter/2), Mm(y - diameter/2),
                                 Mm(diameter), Mm(diameter))
    _clean_shape(dot)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    return diameter


# ============================================================================
# 10. DividerLine — horizontal separator line
# ============================================================================

def DividerLine(slide, x, y, w, color=None, thickness=0.5):
    """Horizontal separator line between content sections.

    Args:
        slide: pptx slide object
        x, y: position in mm
        w: line width in mm
        color: line color (default CATHAY_LTGREY)
        thickness: line thickness in mm (default 0.5)

    Returns:
        bottom_y_mm (float)
    """
    color = color or CATHAY_LTGREY
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Mm(x), Mm(y), Mm(w), Mm(thickness))
    _clean_shape(line)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return y + thickness


# ============================================================================
# 6. Legacy aliases (backward compatibility)
# ============================================================================

def add_callout_box(slide, x_mm, y_mm, w_mm, h_mm, value, label,
                    bg_color=None, text_color=None):
    """Legacy KPI callout box. Prefer KpiStrip for rows of metrics."""
    bg_color = bg_color or CATHAY_RED
    text_color = text_color or CATHAY_WHITE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg_color
    sh.line.fill.background()
    tf = sh.text_frame
    setup_text_frame(tf)
    pv = tf.paragraphs[0]
    pv.alignment = PP_ALIGN.CENTER
    rv = pv.add_run()
    set_run_font(rv, value, size_pt=KPI_VALUE_PT, bold=True, color_rgb=text_color)
    pl = tf.add_paragraph()
    pl.alignment = PP_ALIGN.CENTER
    rl = pl.add_run()
    set_run_font(rl, label, size_pt=KPI_LABEL_PT, color_rgb=text_color)
    return sh


def add_flow_box(slide, x_mm, y_mm, w_mm, h_mm, text, bg_color=None,
                 text_color=None, font_size=9):
    """Legacy flow box. Prefer HeaderBar for labeled sections."""
    bg_color = bg_color or CATHAY_RED
    text_color = text_color or CATHAY_WHITE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg_color
    sh.line.fill.background()
    tf = sh.text_frame
    setup_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_mixed_text(p, text, size_pt=font_size, bold=True, color_rgb=text_color)
    return sh


def add_color_block(slide, x_mm, y_mm, w_mm, h_mm, color=None):
    """Solid color rectangle. Use HeaderBar for labeled blocks."""
    color = color or CATHAY_RED
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_kpi_row(slide, kpis, y_mm=None):
    """Legacy KPI row. Alias for KpiStrip at default position X1, y_mm."""
    y_mm = y_mm or CT
    return KpiStrip(slide, X1, y_mm, kpis)


def add_arrow(slide, x_mm, y_mm, w_mm=8, h_mm=6, color=None):
    """Right-pointing arrow shape."""
    color = color or CATHAY_GOLD
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_down_arrow(slide, x_mm, y_mm, w_mm=6, h_mm=8, color=None):
    """Down-pointing arrow shape."""
    color = color or CATHAY_GOLD
    sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(sh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_progress_bar(slide, x_mm, y_mm, w_mm, h_mm=4, pct=1.0,
                     fill_color=None, bg_color=None):
    """Horizontal progress bar (background + filled portion)."""
    fill_color = fill_color or CATHAY_RED
    bg_color = bg_color or CATHAY_LTGREY

    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    _clean_shape(bg)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    fill_w = max(w_mm * pct, 2)
    fg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Mm(x_mm), Mm(y_mm), Mm(fill_w), Mm(h_mm))
    _clean_shape(fg)
    fg.fill.solid()
    fg.fill.fore_color.rgb = fill_color
    fg.line.fill.background()
    return bg, fg


def add_section_marker(slide, x_mm, y_mm, icon_type=None):
    """Place a small colored shape (4x4mm) as a section visual marker."""
    from constants import ICON_INSIGHT
    icon_type = icon_type or ICON_INSIGHT
    shape_enum, color_hex = icon_type
    marker = slide.shapes.add_shape(shape_enum, Mm(x_mm), Mm(y_mm), Mm(4), Mm(4))
    _clean_shape(marker)
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    marker.line.fill.background()
    return marker


def auto_assign_icons(items):
    """Auto-assign icons to level-0 items based on keyword matching."""
    from constants import _ICON_KEYWORD_MAP, ICON_INSIGHT
    icons = {}
    for text, level in items:
        if level != 0:
            continue
        for icon_type, keywords in _ICON_KEYWORD_MAP.items():
            if any(kw in text for kw in keywords):
                icons[text] = icon_type
                break
        else:
            icons[text] = ICON_INSIGHT
    return icons


__all__ = [
    # New Building Blocks
    "HeaderBar",
    "ContentPanel",
    "KpiStrip",
    "Card",
    "MetricRow",
    "SectionBlock",
    "IconBadge",
    "AccentLine",
    "BulletDot",
    "DividerLine",
    # Legacy aliases (backward compat)
    "add_callout_box",
    "add_flow_box",
    "add_color_block",
    "add_kpi_row",
    "add_arrow",
    "add_down_arrow",
    "add_progress_bar",
    "add_section_marker",
    "auto_assign_icons",
]
