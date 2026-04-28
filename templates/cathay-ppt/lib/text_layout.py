"""
Cathay PPT Template — Text Layout Engine
==========================================
setup_text_frame, format_paragraph, set_square_bullet, add_bullet_content,
add_multi_text, calc_text_height, calc_textframe_height, smart_textbox.

Usage:
    from text_layout import setup_text_frame, smart_textbox, add_bullet_content
"""

import math
import re

from pptx.util import Pt, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree

from constants import (
    MARGIN_ALL, DEFAULT_FONT_SIZE, INDENT_LEFT,
    SPACING_BEFORE, SPACING_AFTER, LINE_SPACING_PCT,
    CL, CT, CB, CW,
    CATHAY_RED, CATHAY_BLACK, CATHAY_GREY, CATHAY_WHITE,
)
from fonts import set_run_font, get_char_width


# ============================================================================
# 1. TEXT FRAME SETUP
# ============================================================================

def setup_text_frame(tf, word_wrap=True):
    """Apply standard Cathay text frame settings.
    - 0.2cm margins all sides
    - No auto-shrink (MSO_AUTO_SIZE.NONE)
    """
    tf.word_wrap = word_wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left   = MARGIN_ALL
    tf.margin_right  = MARGIN_ALL
    tf.margin_top    = MARGIN_ALL
    tf.margin_bottom = MARGIN_ALL


# ============================================================================
# 2. PARAGRAPH FORMATTING
# ============================================================================

def format_paragraph(para, indent_left=True, is_bullet=False):
    """Apply standard Cathay paragraph formatting.
    - Spacing before: 4pt, after: 0pt
    - Line spacing: 1.2x (120%)
    - Left indent: 0.5cm (for bulleted paragraphs)
    """
    pPr = para._p.get_or_add_pPr()

    # Spacing before 4pt, after 0pt
    spcBef = pPr.find(qn('a:spcBef'))
    if spcBef is None:
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
    else:
        spcBef.clear()
    spcPts_bef = etree.SubElement(spcBef, qn('a:spcPts'))
    spcPts_bef.set('val', '400')

    spcAft = pPr.find(qn('a:spcAft'))
    if spcAft is None:
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
    else:
        spcAft.clear()
    spcPts_aft = etree.SubElement(spcAft, qn('a:spcPts'))
    spcPts_aft.set('val', '0')

    # Line spacing 1.2x (120%)
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is None:
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
    else:
        lnSpc.clear()
    spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
    spcPct.set('val', str(LINE_SPACING_PCT))

    # Left indent for bulleted paragraphs
    if indent_left and is_bullet:
        pPr.set('indent', str(-Mm(3)))
        pPr.set('marL', str(INDENT_LEFT))


# ============================================================================
# 3. BULLETS
# ============================================================================

def set_square_bullet(para, color='000000'):
    """Set filled square bullet via PPT buChar XML at 70% size."""
    pPr = para._p.get_or_add_pPr()
    for tag in ('a:buNone', 'a:buChar', 'a:buSzPct', 'a:buClr', 'a:buAutoNum', 'a:buFont'):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    buFont = etree.SubElement(pPr, qn('a:buFont'))
    buFont.set('typeface', 'Calibri')
    buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
    buSzPct.set('val', '70000')
    buClr = etree.SubElement(pPr, qn('a:buClr'))
    srgb = etree.SubElement(buClr, qn('a:srgbClr'))
    srgb.set('val', color)
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', '\u25A0')


# ============================================================================
# 4. BULLETED CONTENT (3-level hierarchy)
# ============================================================================

def add_bullet_content(tf, items, size_pt=None, color_rgb=None):
    """Add bulleted content to a text frame.

    items: list of (text, level) tuples.
      level 0 = section header (bold red, no bullet, no indent)
      level 1 = bulleted item (filled square, 0.5cm indent)
      level 2+ = sub-item (smaller, grey, 1.0cm indent)
    """
    size_pt = size_pt or DEFAULT_FONT_SIZE
    color = color_rgb or CATHAY_BLACK
    for i, (text, level) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        format_paragraph(p, indent_left=True, is_bullet=(level >= 1))

        run = p.add_run()

        if level == 0:
            set_run_font(run, text, size_pt=size_pt + 2, bold=True, color_rgb=CATHAY_RED)
        elif level == 1:
            set_run_font(run, text, size_pt=size_pt, color_rgb=color)
            set_square_bullet(p)
        elif level >= 2:
            set_run_font(run, text, size_pt=max(size_pt - 1, 8), color_rgb=RGBColor(0x80, 0x80, 0x80))
            set_square_bullet(p, color='808080')
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(Mm(10)))
            pPr.set('indent', str(-Mm(3)))


# ============================================================================
# 5. MULTI-TEXT (flexible per-paragraph formatting)
# ============================================================================

def add_multi_text(slide, x_mm, y_mm, w_mm, h_mm, segments, align=PP_ALIGN.LEFT, fill_rgb=None):
    """Add a textbox with multiple formatted paragraphs.

    More flexible than add_bullet_content — each segment has its own formatting.

    Args:
        segments: list of (text, kwargs_dict) tuples.
            kwargs: size (pt), bold, color (RGBColor), italic, space_before (pt),
                    space_after (pt), line_spacing (float, e.g. 1.1)
        fill_rgb: optional RGBColor for textbox background

    Returns:
        textbox shape
    """
    txBox = slide.shapes.add_textbox(Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    tf = txBox.text_frame
    setup_text_frame(tf)

    for idx, (text, kw) in enumerate(segments):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = kw.get('align', align)

        if 'space_before' in kw:
            p.space_before = Pt(kw['space_before'])
        if 'space_after' in kw:
            p.space_after = Pt(kw['space_after'])
        if 'line_spacing' in kw:
            p.line_spacing = kw['line_spacing']

        run = p.add_run()
        set_run_font(
            run, text,
            size_pt=kw.get('size', DEFAULT_FONT_SIZE),
            bold=kw.get('bold', False),
            color_rgb=kw.get('color', CATHAY_BLACK),
        )
        if kw.get('italic', False):
            run.font.italic = True

    if fill_rgb:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill_rgb
    else:
        txBox.fill.background()
    txBox.line.fill.background()

    return txBox


# ============================================================================
# 6. TEXT HEIGHT CALCULATOR
# ============================================================================

def calc_text_height(text_or_paragraphs, box_width_mm, font_pt=10,
                     line_spacing=1.2, margin_mm=4):
    """Calculate rendered text height in mm.

    Args:
        text_or_paragraphs: str or [(text, font_pt, indent_mm), ...]
        box_width_mm: textbox width (mm)
        font_pt: default font size
        line_spacing: line spacing multiplier
        margin_mm: total top+bottom margins (mm)

    Returns:
        float: estimated rendered height (mm)
    """
    usable_w = box_width_mm - margin_mm
    if usable_w <= 0:
        usable_w = 5

    if isinstance(text_or_paragraphs, str):
        paragraphs = [(text_or_paragraphs, font_pt, 0)]
    else:
        paragraphs = text_or_paragraphs

    total_h = margin_mm / 2

    for i, (text, p_font, indent_mm) in enumerate(paragraphs):
        if not text.strip():
            total_h += 1.5
            continue

        effective_w = usable_w - indent_mm
        if effective_w <= 0:
            effective_w = usable_w

        segments = re.findall(
            r'[\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+|[^\u4e00-\u9fff\u3000-\u303f\uff00-\uffef]+',
            text
        )

        total_text_width = 0
        for seg in segments:
            is_cjk = any('\u4e00' <= c <= '\u9fff' for c in seg)
            char_w = get_char_width(p_font, is_cjk)
            total_text_width += len(seg) * char_w

        n_lines = max(1, math.ceil(total_text_width / effective_w))
        line_h = p_font * 0.3528 * line_spacing

        if i > 0:
            total_h += 1.0

        total_h += n_lines * line_h

    total_h += margin_mm / 2
    return total_h


def calc_textframe_height(text_frame, box_width_mm):
    """Calculate rendered height from a python-pptx TextFrame object."""
    paragraph_data = []
    for p in text_frame.paragraphs:
        p_font = 10
        for r in p.runs:
            if r.font.size:
                p_font = r.font.size / 12700
                break

        indent_mm = 0
        pPr = p._p.get_or_add_pPr()
        marL = pPr.get('marL')
        if marL:
            indent_mm = int(marL) / 36000

        paragraph_data.append((p.text, p_font, indent_mm))

    line_sp = 1.2
    if text_frame.paragraphs:
        pPr_el = text_frame.paragraphs[0]._p.find(qn('a:lnSpc'))
        if pPr_el is not None:
            spcPct = pPr_el.find(qn('a:spcPct'))
            if spcPct is not None:
                line_sp = int(spcPct.get('val', '120000')) / 100000

    m_tb = ((text_frame.margin_top or 36000) / 36000 +
            (text_frame.margin_bottom or 36000) / 36000)

    return calc_text_height(paragraph_data, box_width_mm, line_spacing=line_sp, margin_mm=m_tb)


# ============================================================================
# 7. SMART TEXTBOX (auto-fit)
# ============================================================================

def smart_textbox(slide, x_mm, y_mm, w_mm, items, max_bottom_mm=180,
                   start_font=10, min_font=8, line_spacing=1.2):
    """Create auto-fit textbox: calculate height, reduce font if overflow.

    Args:
        slide: pptx slide object
        x_mm, y_mm, w_mm: position and width in mm
        items: [(text, level), ...] for add_bullet_content
        max_bottom_mm: bottom limit (mm)
        start_font: initial font size (pt)
        min_font: minimum font size (pt)

    Returns:
        (shape, text_frame, actual_font_pt)
    """
    max_h = max_bottom_mm - y_mm

    chosen_font = start_font
    est_h = max_h + 1
    for try_font in [start_font, start_font - 0.5, start_font - 1,
                     start_font - 1.5, start_font - 2, min_font]:
        if try_font < min_font:
            try_font = min_font

        paras = []
        for text, level in items:
            indent = 5 if level == 1 else (10 if level >= 2 else 0)
            f = try_font + 1 if level == 0 else (try_font - 1 if level >= 2 else try_font)
            paras.append((text, f, indent))

        est_h = calc_text_height(paras, w_mm, line_spacing=line_spacing)

        if est_h <= max_h:
            chosen_font = try_font
            break
        chosen_font = try_font

    actual_h = min(max_h, est_h + 3)
    txBox = slide.shapes.add_textbox(Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(actual_h))
    tf = txBox.text_frame
    setup_text_frame(tf)
    add_bullet_content(tf, items, size_pt=chosen_font)

    return txBox, tf, chosen_font


__all__ = [
    "setup_text_frame",
    "format_paragraph",
    "set_square_bullet",
    "add_bullet_content",
    "add_multi_text",
    "calc_text_height",
    "calc_textframe_height",
    "smart_textbox",
]
