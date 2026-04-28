"""
QC Automation Module for Cathay PPT Template.

Provides:
  - full_qc_pipeline(): validate shapes -> export PDF -> convert PNG -> report
  - check_guard_rails(): 8 production guard rails (McKinsey-adapted for Cathay 4:3)
  - autofix_pipeline(): 4-stage auto-fix chain (remove -> compress -> restructure -> font)
  - update_slide_in_deck(): replace a single slide without full regeneration
  - batch_validate(): validate all slide files in a directory
  - auto_fix_all(): run validate_and_fix on every slide file
"""

import os
import glob
import copy
import io
import math
import re
import subprocess

from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree

# ── Import shared modules ──
from constants import (
    CL, CT, CW, CH,
    CONTENT_BOTTOM_MM,
    SOURCE_Y_MM,
    SOURCE_FONT_PT,
    MIN_SOURCE_FONT_PT,
    MIN_BODY_FONT_PT,
)
from fonts import get_char_width
from text_layout import calc_textframe_height, calc_text_height
from validation import validate_and_fix, save_with_validation
from merge import merge_slides

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
CONTENT_ZONE_BOTTOM_MM = 181   # shapes must not exceed this (except footer)

# Guard rail constants (Cathay 4:3 slide = 254mm x 190.5mm)
SLIDE_RIGHT_MARGIN_MM = 244    # max right edge for content
FOOTER_ZONE_TOP_MM = 181       # footer zone starts here
COLLISION_TOLERANCE_MM = 0.8   # min gap between non-contained shapes
BOTTOM_WHITESPACE_MAX_MM = 8   # max gap between lowest content and footer zone
PEER_FONT_Y_TOLERANCE_MM = 2   # shapes at same Y should match fonts

# CJK character density limits: {box_height_mm: max_chars}
CJK_DENSITY_LIMITS = {
    5: 15, 10: 40, 15: 80, 20: 130, 25: 190,
    30: 260, 40: 450, 50: 700,
}

AUTOFIX_TITLE_MIN_PT = 18
AUTOFIX_BODY_MIN_PT = 9
AUTOFIX_FOOTER_MIN_PT = 7


# ═══════════════════════════════════════════════════════════════════════════
# 1a. Guard Rail Helpers (private)
# ═══════════════════════════════════════════════════════════════════════════

def _is_footer_shape(shape):
    """Detect footer / source / page-number shapes."""
    if not shape.has_text_frame:
        return False
    txt = shape.text_frame.text.lower()
    if "source:" in txt or "source\uff1a" in txt:
        return True
    if len(txt) < 10 and "/" in txt:
        return True
    top_mm = shape.top / 36000
    if top_mm > 178:
        if len(txt.strip()) < 80:
            return True
    return False


def _get_layout_signature(slide):
    """Return a rough grid-pattern key based on shape count and quantized positions."""
    buckets = set()
    count = 0
    for shape in slide.shapes:
        w_mm = shape.width / 36000
        h_mm = shape.height / 36000
        if w_mm < 1 and h_mm < 1:
            continue
        if _is_footer_shape(shape):
            continue
        count += 1
        col = min(int((shape.left / 36000) / 64), 3)
        row = min(int((shape.top / 36000) / 48), 3)
        buckets.add((row, col))
    return f"{count}:{','.join(sorted(f'{r}{c}' for r, c in buckets))}"


def _check_char_density(shape, height_mm, slide_num, violations):
    """Check CJK text against density limits per box height."""
    text = shape.text_frame.text
    cjk_count = sum(1 for c in text if "\u4e00" <= c <= "\u9fff"
                    or "\u3000" <= c <= "\u303f"
                    or "\uff00" <= c <= "\uffef")
    if cjk_count == 0:
        return

    max_allowed = None
    for h_threshold in sorted(CJK_DENSITY_LIMITS.keys()):
        if h_threshold <= height_mm:
            max_allowed = CJK_DENSITY_LIMITS[h_threshold]
        else:
            break

    if max_allowed is None:
        max_allowed = CJK_DENSITY_LIMITS[5]

    if cjk_count > max_allowed:
        violations.append(
            f"S{slide_num} GUARD8_DENSITY: '{shape.name}' {cjk_count} CJK chars "
            f"in {height_mm:.0f}mm box (limit ~{max_allowed})"
        )


def _check_peer_fonts(slide, slide_num, violations):
    """Shapes at the same Y (within tolerance) should have the same font size."""
    font_records = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if _is_footer_shape(shape):
            continue
        top_mm = shape.top / 36000
        sizes = []
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    sizes.append(r.font.size / 12700)
        if not sizes:
            continue
        dominant_pt = max(set(sizes), key=sizes.count)
        font_records.append((top_mm, dominant_pt, shape.name))

    used = set()
    for i in range(len(font_records)):
        if i in used:
            continue
        group = [font_records[i]]
        used.add(i)
        for j in range(i + 1, len(font_records)):
            if j in used:
                continue
            if abs(font_records[j][0] - font_records[i][0]) <= PEER_FONT_Y_TOLERANCE_MM:
                group.append(font_records[j])
                used.add(j)

        if len(group) < 2:
            continue

        sizes_in_group = set(round(g[1], 1) for g in group)
        if len(sizes_in_group) > 1:
            size_str = ", ".join(f"{g[2]}={g[1]:.0f}pt" for g in group)
            violations.append(
                f"S{slide_num} GUARD5_PEER_FONT: shapes at Y~{group[0][0]:.0f}mm "
                f"have mismatched fonts: {size_str}"
            )


def _check_collisions(slide, slide_num, violations):
    """Detect non-contained shapes within 0.8mm vertical gap of each other."""
    boxes = []
    for shape in slide.shapes:
        w_mm = shape.width / 36000
        h_mm = shape.height / 36000
        if w_mm < 1 and h_mm < 1:
            continue
        if _is_footer_shape(shape):
            continue
        l = shape.left / 36000
        t = shape.top / 36000
        r = l + w_mm
        b = t + h_mm
        boxes.append((shape.name, l, t, r, b))

    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            n1, l1, t1, r1, b1 = boxes[i]
            n2, l2, t2, r2, b2 = boxes[j]

            if r1 <= l2 or r2 <= l1:
                continue

            if _contains(l1, t1, r1, b1, l2, t2, r2, b2):
                continue
            if _contains(l2, t2, r2, b2, l1, t1, r1, b1):
                continue

            vert_gap = max(t2 - b1, t1 - b2)
            if 0 < vert_gap < COLLISION_TOLERANCE_MM:
                violations.append(
                    f"S{slide_num} GUARD6_COLLISION: '{n1}' and '{n2}' "
                    f"only {vert_gap:.1f}mm apart vertically"
                )


def _check_bottom_whitespace(slide, slide_num, violations):
    """Flag slides where content doesn't fill close to the footer zone."""
    max_bottom = 0
    has_content = False
    for shape in slide.shapes:
        if _is_footer_shape(shape):
            continue
        w_mm = shape.width / 36000
        h_mm = shape.height / 36000
        if w_mm < 1 and h_mm < 1:
            continue
        has_content = True
        bottom = shape.top / 36000 + h_mm
        if bottom > max_bottom:
            max_bottom = bottom

    if has_content and max_bottom > 0:
        gap = FOOTER_ZONE_TOP_MM - max_bottom
        if gap > BOTTOM_WHITESPACE_MAX_MM:
            violations.append(
                f"S{slide_num} GUARD3_WHITESPACE: lowest content at {max_bottom:.0f}mm, "
                f"{gap:.0f}mm gap to footer zone (>{BOTTOM_WHITESPACE_MAX_MM}mm)"
            )


def _check_layout_variety(layouts, violations):
    """Check layout variety: 5+ unique per 25 slides, no 3 consecutive identical."""
    if len(layouts) >= 25:
        unique = len(set(layouts[:25]))
        if unique < 5:
            violations.append(
                f"GUARD7_VARIETY: only {unique} unique layouts in first 25 slides (need 5+)"
            )

    for i in range(len(layouts) - 2):
        if layouts[i] == layouts[i + 1] == layouts[i + 2]:
            violations.append(
                f"GUARD7_VARIETY: slides {i+1}-{i+3} have identical layout '{layouts[i]}'"
            )
            break


# ═══════════════════════════════════════════════════════════════════════════
# 1b. Guard Rail Check (8 rules)
# ═══════════════════════════════════════════════════════════════════════════

def check_guard_rails(prs):
    """Check 8 production guard rails (McKinsey-adapted for Cathay 4:3)."""
    violations = []
    slide_layouts_used = []

    for slide_idx, slide in enumerate(prs.slides):
        sn = slide_idx + 1
        layout_key = _get_layout_signature(slide)
        slide_layouts_used.append(layout_key)

        for shape in slide.shapes:
            top_mm = shape.top / 36000
            left_mm = shape.left / 36000
            width_mm = shape.width / 36000
            height_mm = shape.height / 36000
            bottom_mm = top_mm + height_mm
            right_mm = left_mm + width_mm

            if width_mm < 1 and height_mm < 1:
                continue

            is_footer = _is_footer_shape(shape)

            if not is_footer and bottom_mm > (FOOTER_ZONE_TOP_MM - 1):
                violations.append(
                    f"S{sn} GUARD1_GAP: '{shape.name}' bottom={bottom_mm:.1f}mm, "
                    f"<1mm from footer zone"
                )

            if right_mm > SLIDE_RIGHT_MARGIN_MM:
                violations.append(
                    f"S{sn} GUARD2_MARGIN: '{shape.name}' right={right_mm:.1f}mm "
                    f"> {SLIDE_RIGHT_MARGIN_MM}mm"
                )

            if shape.has_text_frame:
                _check_char_density(shape, height_mm, sn, violations)

        _check_peer_fonts(slide, sn, violations)
        _check_collisions(slide, sn, violations)
        _check_bottom_whitespace(slide, sn, violations)

    _check_layout_variety(slide_layouts_used, violations)

    return violations


# ═══════════════════════════════════════════════════════════════════════════
# 1c. AutoFix Priority Chain (4 stages)
# ═══════════════════════════════════════════════════════════════════════════

def autofix_pipeline(pptx_path):
    """Four-stage auto-fix pipeline (McKinsey pattern, Cathay-adapted)."""
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    prs = Presentation(pptx_path)
    result = {"stage1": [], "stage2": [], "stage3": [], "stage4": []}

    for slide_idx, slide in enumerate(prs.slides):
        sn = slide_idx + 1
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            seen_texts = []
            paras_to_remove = []
            for p_idx, para in enumerate(tf.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                if text in seen_texts:
                    paras_to_remove.append(p_idx)
                else:
                    seen_texts.append(text)

            if paras_to_remove:
                for p_idx in reversed(paras_to_remove):
                    p_elem = tf.paragraphs[p_idx]._p
                    p_elem.getparent().remove(p_elem)
                result["stage1"].append(
                    f"S{sn} '{shape.name}': removed {len(paras_to_remove)} duplicate paragraph(s)"
                )

    for slide_idx, slide in enumerate(prs.slides):
        sn = slide_idx + 1
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    original = run.text
                    compressed = re.sub(r"  +", " ", original)
                    compressed = compressed.strip()
                    if compressed != original:
                        run.text = compressed
                        if not result["stage2"] or result["stage2"][-1].split(":")[0] != f"S{sn} '{shape.name}'":
                            result["stage2"].append(
                                f"S{sn} '{shape.name}': compressed whitespace"
                            )

    for slide_idx, slide in enumerate(prs.slides):
        sn = slide_idx + 1
        for shape in slide.shapes:
            if _is_footer_shape(shape):
                continue

            top_mm = shape.top / 36000
            left_mm = shape.left / 36000
            width_mm = shape.width / 36000
            height_mm = shape.height / 36000
            bottom_mm = top_mm + height_mm
            right_mm = left_mm + width_mm

            if bottom_mm > CONTENT_ZONE_BOTTOM_MM:
                new_height_mm = CONTENT_ZONE_BOTTOM_MM - top_mm
                if new_height_mm > 0:
                    shape.height = int(new_height_mm * 36000)
                    result["stage3"].append(
                        f"S{sn} '{shape.name}': capped height from {height_mm:.0f}mm "
                        f"to {new_height_mm:.0f}mm (content zone)"
                    )
                else:
                    shape.top = int((CONTENT_ZONE_BOTTOM_MM - height_mm) * 36000)
                    result["stage3"].append(
                        f"S{sn} '{shape.name}': repositioned above content zone bottom"
                    )

            if right_mm > SLIDE_RIGHT_MARGIN_MM:
                overflow = right_mm - SLIDE_RIGHT_MARGIN_MM
                new_left_mm = left_mm - overflow
                if new_left_mm >= CL:
                    shape.left = int(new_left_mm * 36000)
                    result["stage3"].append(
                        f"S{sn} '{shape.name}': shifted left by {overflow:.0f}mm (right margin)"
                    )
                else:
                    new_width_mm = SLIDE_RIGHT_MARGIN_MM - left_mm
                    if new_width_mm > 5:
                        shape.width = int(new_width_mm * 36000)
                        result["stage3"].append(
                            f"S{sn} '{shape.name}': trimmed width to {new_width_mm:.0f}mm (right margin)"
                        )

    for slide_idx, slide in enumerate(prs.slides):
        sn = slide_idx + 1
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            top_mm = shape.top / 36000
            height_mm = shape.height / 36000
            width_mm = shape.width / 36000

            if height_mm < 5:
                continue

            is_footer = _is_footer_shape(shape)
            is_title = top_mm < 30

            if is_footer:
                min_pt = AUTOFIX_FOOTER_MIN_PT
            elif is_title:
                min_pt = AUTOFIX_TITLE_MIN_PT
            else:
                min_pt = AUTOFIX_BODY_MIN_PT

            est_h = calc_textframe_height(shape.text_frame, width_mm)
            if est_h <= height_mm * 1.05:
                continue

            adjusted = False
            for _ in range(8):
                any_reduced = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            current_pt = run.font.size / 12700
                            new_pt = current_pt - 0.5
                            if new_pt >= min_pt:
                                run.font.size = int(new_pt * 12700)
                                any_reduced = True
                if not any_reduced:
                    break
                est_h = calc_textframe_height(shape.text_frame, width_mm)
                if est_h <= height_mm * 1.05:
                    adjusted = True
                    break

            if adjusted:
                result["stage4"].append(
                    f"S{sn} '{shape.name}': reduced fonts to fit (floor: {min_pt}pt)"
                )

    total_fixes = sum(len(v) for v in result.values())
    if total_fixes > 0:
        output_path = pptx_path.replace(".pptx", "_autofix.pptx")
        prs.save(output_path)
        result["output_path"] = output_path
    else:
        result["output_path"] = pptx_path

    print(f"AutoFix pipeline: {total_fixes} total fixes across 4 stages")
    for stage_name in ("stage1", "stage2", "stage3", "stage4"):
        fixes = result[stage_name]
        if fixes:
            print(f"  {stage_name}: {len(fixes)} fix(es)")
            for f in fixes[:3]:
                print(f"    - {f}")
            if len(fixes) > 3:
                print(f"    ... and {len(fixes) - 3} more")

    return result


# ═══════════════════════════════════════════════════════════════════════════
# 2. PDF Visual QC Pipeline
# ═══════════════════════════════════════════════════════════════════════════

def full_qc_pipeline(pptx_path, output_dir=None):
    """Complete QC: autofix -> validate -> guard rails -> PDF -> PNG -> report."""
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    if output_dir is None:
        output_dir = os.path.dirname(os.path.abspath(pptx_path))
    os.makedirs(output_dir, exist_ok=True)

    autofix_result = autofix_pipeline(pptx_path)
    working_path = autofix_result.get("output_path", pptx_path)

    prs = Presentation(working_path)
    fixes = validate_and_fix(prs)
    if fixes:
        fixed_path = pptx_path.replace(".pptx", "_qc_fixed.pptx")
        prs.save(fixed_path)
        working_path = fixed_path
    else:
        fixed_path = None

    prs_for_guard = Presentation(working_path)
    guard_violations = check_guard_rails(prs_for_guard)

    pdf_path = _export_pdf(working_path, output_dir)

    png_dir = os.path.join(output_dir, "qc_pages")
    if pdf_path:
        _pdf_to_pngs(pdf_path, png_dir)

    prs = Presentation(working_path)
    report = {}

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        issues = []

        has_source_footer = False
        shapes_for_overlap = []

        for shape in slide.shapes:
            top_mm = shape.top / 36000
            left_mm = shape.left / 36000
            width_mm = shape.width / 36000
            height_mm = shape.height / 36000
            bottom_mm = top_mm + height_mm
            right_mm = left_mm + width_mm

            if width_mm < 0.5 or height_mm < 0.5:
                continue

            is_footer = False
            if shape.has_text_frame:
                txt = shape.text_frame.text.lower()
                if "source:" in txt or "source\uff1a" in txt:
                    has_source_footer = True
                    is_footer = True
                if len(txt) < 10 and "/" in txt:
                    is_footer = True

            if not is_footer and bottom_mm > CONTENT_ZONE_BOTTOM_MM:
                issues.append(
                    f"OVERFLOW: '{shape.name}' bottom={bottom_mm:.1f}mm > {CONTENT_ZONE_BOTTOM_MM}mm"
                )

            if shape.has_text_frame and height_mm >= 5:
                est_h = calc_textframe_height(shape.text_frame, width_mm)
                if est_h > height_mm * 1.1:
                    issues.append(
                        f"TEXT_OVERFLOW: '{shape.name}' est={est_h:.1f}mm > box={height_mm:.1f}mm"
                    )

            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            pt = r.font.size / 12700
                            if is_footer and pt < 7:
                                issues.append(
                                    f"FONT_TOO_SMALL: '{shape.name}' footer font {pt:.1f}pt < 7pt"
                                )
                            elif not is_footer and pt < 9:
                                issues.append(
                                    f"FONT_TOO_SMALL: '{shape.name}' content font {pt:.1f}pt < 9pt"
                                )

            if shape.has_table:
                table = shape.table
                for ri in range(len(table.rows)):
                    row_h_mm = table.rows[ri].height / 36000 if table.rows[ri].height else 0
                    for ci in range(len(table.columns)):
                        cell = table.cell(ri, ci)
                        cell_text = cell.text_frame.text
                        if not cell_text.strip():
                            continue
                        col_w_mm = table.columns[ci].width / 36000 if table.columns[ci].width else 30
                        cell_font = 9
                        for cp in cell.text_frame.paragraphs:
                            for cr in cp.runs:
                                if cr.font.size:
                                    cell_font = cr.font.size / 12700
                                    break
                            break
                        has_cjk = any("\u4e00" <= c <= "\u9fff" for c in cell_text)
                        char_w = get_char_width(cell_font, has_cjk)
                        usable = max(col_w_mm - 3, 5)
                        lines = max(1, math.ceil(len(cell_text) * char_w / usable))
                        line_h = cell_font * 0.3528 * 1.2
                        needed = lines * line_h + 3
                        if row_h_mm > 0 and needed > row_h_mm * 1.15:
                            issues.append(
                                f"TABLE_ROW_SHORT: '{shape.name}' row {ri} needs {needed:.1f}mm, has {row_h_mm:.1f}mm"
                            )
                    break

            if not is_footer:
                shapes_for_overlap.append(
                    (shape.name, left_mm, top_mm, right_mm, bottom_mm)
                )

        for i in range(len(shapes_for_overlap)):
            for j in range(i + 1, len(shapes_for_overlap)):
                n1, l1, t1, r1, b1 = shapes_for_overlap[i]
                n2, l2, t2, r2, b2 = shapes_for_overlap[j]
                if l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2:
                    if _contains(l1, t1, r1, b1, l2, t2, r2, b2):
                        continue
                    if _contains(l2, t2, r2, b2, l1, t1, r1, b1):
                        continue
                    issues.append(f"OVERLAP: '{n1}' and '{n2}'")

        if not has_source_footer:
            issues.append("MISSING_SOURCE: no 'Source:' footer found")

        if issues:
            report[slide_num] = issues

    report["guard_rails"] = guard_violations
    report["autofix"] = autofix_result

    slide_issues = {k: v for k, v in report.items() if isinstance(k, int)}
    total_issues = sum(len(v) for v in slide_issues.values())
    total_autofix = sum(len(v) for k, v in autofix_result.items() if k.startswith("stage"))
    print(f"QC Pipeline complete: {len(prs.slides)} slides, {total_issues} issues, "
          f"{len(guard_violations)} guard-rail violations, {total_autofix} auto-fixes applied")
    if fixes:
        print(f"  validate_and_fix: {len(fixes)} fixes (saved to {fixed_path})")
    if guard_violations:
        print(f"  Guard rail violations ({len(guard_violations)}):")
        for v in guard_violations[:8]:
            print(f"    - {v}")
        if len(guard_violations) > 8:
            print(f"    ... and {len(guard_violations) - 8} more")
    for sn, iss in sorted(slide_issues.items()):
        print(f"  Slide {sn}: {len(iss)} issues")
        for i in iss[:5]:
            print(f"    - {i}")
        if len(iss) > 5:
            print(f"    ... and {len(iss) - 5} more")

    return report


def _contains(l1, t1, r1, b1, l2, t2, r2, b2):
    return l1 <= l2 and t1 <= t2 and r1 >= r2 and b1 >= b2


def _export_pdf(pptx_path, output_dir):
    """Export PPTX to PDF using LibreOffice headless."""
    try:
        soffice = None
        candidates = [
            "soffice", "/usr/bin/soffice", "/usr/local/bin/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ]
        for c in candidates:
            try:
                subprocess.run([c, "--version"], capture_output=True, timeout=5)
                soffice = c
                break
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue

        if soffice is None:
            print("  [QC] LibreOffice not found — skipping PDF export")
            return None

        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
            capture_output=True, timeout=120,
        )
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(output_dir, base + ".pdf")
        if os.path.exists(pdf_path):
            print(f"  [QC] PDF exported: {pdf_path}")
            return pdf_path
        else:
            print("  [QC] PDF export failed — file not created")
            return None
    except Exception as e:
        print(f"  [QC] PDF export error: {e}")
        return None


def _pdf_to_pngs(pdf_path, output_dir):
    """Convert PDF pages to PNG using pdftoppm."""
    os.makedirs(output_dir, exist_ok=True)
    try:
        subprocess.run(
            ["pdftoppm", "-png", "-r", "150", pdf_path, os.path.join(output_dir, "page")],
            capture_output=True, timeout=120,
        )
        pngs = sorted(glob.glob(os.path.join(output_dir, "page-*.png")))
        if pngs:
            print(f"  [QC] {len(pngs)} PNG pages generated in {output_dir}")
        else:
            print("  [QC] pdftoppm produced no PNGs")
    except FileNotFoundError:
        print("  [QC] pdftoppm not found — skipping PNG conversion")
    except Exception as e:
        print(f"  [QC] PNG conversion error: {e}")


# ═══════════════════════════════════════════════════════════════════════════
# 3. Incremental Slide Update
# ═══════════════════════════════════════════════════════════════════════════

def update_slide_in_deck(deck_path, slide_index, new_slide_path, output_path=None):
    """Replace a single slide in an existing deck without regenerating everything."""
    if output_path is None:
        output_path = deck_path

    deck = Presentation(deck_path)
    total = len(deck.slides)
    if slide_index < 0 or slide_index >= total:
        raise IndexError(f"slide_index {slide_index} out of range (deck has {total} slides)")

    rId = deck.slides._sldIdLst[slide_index].rId
    deck.part.drop_rel(rId)
    del deck.slides._sldIdLst[slide_index]

    src_prs = Presentation(new_slide_path)
    if len(src_prs.slides) == 0:
        raise ValueError(f"No slides found in {new_slide_path}")
    src_slide = src_prs.slides[0]

    layout_name = src_slide.slide_layout.name
    target_layout = deck.slide_layouts[4]
    for layout in deck.slide_layouts:
        if layout.name == layout_name:
            target_layout = layout
            break

    new_slide = deck.slides.add_slide(target_layout)

    img_map = {}
    for rel in src_slide.part.rels.values():
        if "image" in str(rel.reltype):
            try:
                img_map[rel.rId] = rel.target_part.blob
            except Exception:
                pass

    rId_remap = {}
    for old_rId, blob in img_map.items():
        image_part, new_rId = new_slide.part.get_or_add_image_part(io.BytesIO(blob))
        rId_remap[old_rId] = new_rId

    for shape in src_slide.shapes:
        el = copy.deepcopy(shape._element)
        for blip in el.findall(".//" + qn("a:blip")):
            old_rId = blip.get(qn("r:embed"))
            if old_rId in rId_remap:
                blip.set(qn("r:embed"), rId_remap[old_rId])
        new_slide.shapes._spTree.append(el)

    sldIdLst = deck.slides._sldIdLst
    new_sldId = sldIdLst[-1]
    sldIdLst.remove(new_sldId)
    sldIdLst.insert(slide_index, new_sldId)

    deck.save(output_path)
    print(f"Slide {slide_index} replaced in {output_path}")
    return output_path


# ═══════════════════════════════════════════════════════════════════════════
# 4. Batch Validate & Auto-Fix
# ═══════════════════════════════════════════════════════════════════════════

def batch_validate(slides_dir, pattern="slide_*.pptx"):
    """Validate all individual slide files in a directory."""
    files = sorted(glob.glob(os.path.join(slides_dir, pattern)))
    if not files:
        print(f"No files matching '{pattern}' in {slides_dir}")
        return {}

    report = {}
    for fpath in files:
        fname = os.path.basename(fpath)
        try:
            prs = Presentation(fpath)
            issues = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    top_mm = shape.top / 36000
                    height_mm = shape.height / 36000
                    width_mm = shape.width / 36000
                    bottom_mm = top_mm + height_mm

                    if width_mm < 0.5 or height_mm < 0.5:
                        continue

                    is_footer = False
                    if shape.has_text_frame:
                        txt = shape.text_frame.text.lower()
                        if "source:" in txt:
                            is_footer = True

                    if not is_footer and bottom_mm > CONTENT_ZONE_BOTTOM_MM:
                        issues.append(f"OVERFLOW: '{shape.name}' bottom={bottom_mm:.1f}mm")

                    if shape.has_text_frame and height_mm >= 5:
                        est_h = calc_textframe_height(shape.text_frame, width_mm)
                        if est_h > height_mm * 1.1:
                            issues.append(
                                f"TEXT_OVERFLOW: '{shape.name}' est={est_h:.1f}mm > box={height_mm:.1f}mm"
                            )

                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                if r.font.size:
                                    pt = r.font.size / 12700
                                    if not is_footer and pt < 9:
                                        issues.append(f"FONT_SMALL: '{shape.name}' {pt:.1f}pt")

            if issues:
                report[fname] = issues
            else:
                report[fname] = []

        except Exception as e:
            report[fname] = [f"ERROR: {e}"]

    total_files = len(files)
    files_with_issues = sum(1 for v in report.values() if v)
    total_issues = sum(len(v) for v in report.values())
    print(f"Batch validate: {total_files} files, {files_with_issues} with issues, {total_issues} total issues")

    return report


def auto_fix_all(slides_dir, pattern="slide_*.pptx"):
    """Run validate_and_fix on every slide file."""
    from validation import validate_and_fix as vf

    files = sorted(glob.glob(os.path.join(slides_dir, pattern)))
    if not files:
        print(f"No files matching '{pattern}' in {slides_dir}")
        return 0

    total_fixes = 0
    for fpath in files:
        fname = os.path.basename(fpath)
        try:
            prs = Presentation(fpath)
            fixes = vf(prs)
            if fixes:
                prs.save(fpath)
                total_fixes += len(fixes)
                print(f"  {fname}: {len(fixes)} fixes applied")
                for f in fixes[:3]:
                    print(f"    - {f}")
                if len(fixes) > 3:
                    print(f"    ... and {len(fixes) - 3} more")
        except Exception as e:
            print(f"  {fname}: ERROR — {e}")

    print(f"Auto-fix complete: {total_fixes} total fixes across {len(files)} files")
    return total_fixes


__all__ = [
    "full_qc_pipeline",
    "check_guard_rails",
    "autofix_pipeline",
    "update_slide_in_deck",
    "batch_validate",
    "auto_fix_all",
]
