"""
Data-Driven Slide Generation for Cathay PPT Template.

Separates content from layout:
  - DataRegistry: single source of truth for all data points
  - Slide specs: declarative dict format for slide content
  - build_deck_from_specs(): specs list -> complete deck
  - render_spec(): render a single spec into a presentation
"""

import os
import json

from pptx import Presentation
from pptx.util import Mm

_LIB_DIR = os.path.dirname(os.path.abspath(__file__))
import sys
if _LIB_DIR not in sys.path:
    sys.path.insert(0, _LIB_DIR)

from constants import (
    TEMPLATE,
    CL, CT, CW, CH, GAP_H, GAP_V,
    CATHAY_RED, CATHAY_GOLD, CATHAY_WHITE, CATHAY_BLACK, CATHAY_GREY,
)
from fonts import set_run_font
from text_layout import setup_text_frame, add_bullet_content
from tables import add_table
from charts import safe_chart_insert
from slides import add_source_footer, set_title_with_conclusion
from elements import add_kpi_row
from safe_layout import safe_textbox
from validation import validate_and_fix
from qc_automation import full_qc_pipeline
from slide_templates import (
    template_kpi_dashboard,
    template_value_chain_flow,
    template_chart_plus_analysis,
    template_comparison_matrix,
    template_two_column_analysis,
    template_sidebar_case_study,
    template_three_column_compare,
    template_stacked_cases,
    template_risk_cards,
    template_action_plan,
    template_donut_chart,
    template_before_after,
    template_funnel,
    template_swot,
    template_waterfall,
    template_stakeholder_map,
    template_timeline,
    template_number_story,
    template_executive_summary,
)


# ═══════════════════════════════════════════════════════════════════════════
# 1. DataRegistry — Single Source of Truth
# ═══════════════════════════════════════════════════════════════════════════

class DataRegistry:
    """Single source of truth for all data points used in a deck.

    Stores (value, source, year) triples. Load from / save to JSON.
    """
    def __init__(self, data_file=None):
        self.data = {}
        self.sources = {}
        self.years = {}
        if data_file:
            self.load(data_file)

    def set(self, key, value, source, year=None):
        self.data[key] = value
        self.sources[key] = source
        if year is not None:
            self.years[key] = year

    def get(self, key, default=None):
        return self.data.get(key, default)

    def get_with_source(self, key):
        if key not in self.data:
            return (None, None)
        return (self.data[key], self.sources.get(key, "unknown"))

    def get_source(self, key):
        return self.sources.get(key)

    def keys(self):
        return list(self.data.keys())

    def __contains__(self, key):
        return key in self.data

    def __len__(self):
        return len(self.data)

    def save(self, path):
        payload = {"data": self.data, "sources": self.sources, "years": self.years}
        with open(path, "w", encoding="utf-8") as f:
            json.dump(payload, f, ensure_ascii=False, indent=2)
        print(f"DataRegistry saved: {len(self.data)} entries -> {path}")

    def load(self, path):
        with open(path, "r", encoding="utf-8") as f:
            payload = json.load(f)
        self.data = payload.get("data", {})
        self.sources = payload.get("sources", {})
        self.years = payload.get("years", {})
        print(f"DataRegistry loaded: {len(self.data)} entries from {path}")

    def validate(self):
        missing = []
        for key in self.data:
            if key not in self.sources or not self.sources[key]:
                missing.append(key)
        if missing:
            print(f"DataRegistry validation: {len(missing)} keys missing sources:")
            for k in missing[:10]:
                print(f"  - {k}")
        else:
            print(f"DataRegistry validation: all {len(self.data)} keys have sources")
        return missing

    def set_many(self, entries):
        for entry in entries:
            if len(entry) == 3:
                self.set(entry[0], entry[1], entry[2])
            elif len(entry) >= 4:
                self.set(entry[0], entry[1], entry[2], entry[3])

    def collect_sources(self, keys):
        sources = set()
        for k in keys:
            s = self.sources.get(k)
            if s:
                sources.add(s)
        return ", ".join(sorted(sources))

    def __repr__(self):
        return f"<DataRegistry: {len(self.data)} entries>"


# ═══════════════════════════════════════════════════════════════════════════
# 2. Template Router
# ═══════════════════════════════════════════════════════════════════════════

TEMPLATE_ROUTER = {
    "kpi_dashboard": template_kpi_dashboard,
    "value_chain_flow": template_value_chain_flow,
    "chart_plus_analysis": template_chart_plus_analysis,
    "comparison_matrix": template_comparison_matrix,
    "two_column_analysis": template_two_column_analysis,
    "sidebar_case_study": template_sidebar_case_study,
    "three_column_compare": template_three_column_compare,
    "stacked_cases": template_stacked_cases,
    "risk_cards": template_risk_cards,
    "action_plan": template_action_plan,
    "donut_chart": template_donut_chart,
    "before_after": template_before_after,
    "funnel": template_funnel,
    "swot": template_swot,
    "waterfall": template_waterfall,
    "stakeholder_map": template_stakeholder_map,
    "timeline": template_timeline,
    "number_story": template_number_story,
    "executive_summary": template_executive_summary,
}


def render_spec(prs, spec, data_registry=None):
    """Render a single slide spec into a presentation."""
    template_name = spec.get("template", "kpi_dashboard")
    renderer = TEMPLATE_ROUTER.get(template_name)

    if renderer is None:
        raise ValueError(
            f"Unknown template '{template_name}'. Available: {list(TEMPLATE_ROUTER.keys())}"
        )

    title = spec.get("title", "")
    subtitle = spec.get("subtitle", "")
    source = spec.get("source", "")
    data = spec.get("data", {})

    if data_registry is not None and not source and "data_keys" in spec:
        source = data_registry.collect_sources(spec["data_keys"])

    try:
        slide = renderer(prs, title, subtitle, **data, source=source)
    except TypeError as e:
        print(f"WARNING: Template '{template_name}' call failed: {e}")
        slide = None

    return slide


# ═══════════════════════════════════════════════════════════════════════════
# 3. Deck Builder
# ═══════════════════════════════════════════════════════════════════════════

def build_deck_from_specs(specs, output_path, template_path=None, data_registry=None,
                          run_qc=True):
    """Build a complete deck from a list of slide specs."""
    template_path = template_path or TEMPLATE
    prs = Presentation(template_path)

    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    for i, spec in enumerate(specs):
        if "page_num" not in spec:
            spec["page_num"] = i + 1
        try:
            render_spec(prs, spec, data_registry=data_registry)
        except Exception as e:
            print(f"Error rendering slide {i + 1} ('{spec.get('template', '?')}'): {e}")
            slide = prs.slides.add_slide(prs.slide_layouts[4])
            for shape in slide.shapes:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                    if shape.placeholder_format.type == 1:
                        shape.text = f"ERROR: Slide {i + 1}"
            txBox = slide.shapes.add_textbox(Mm(CL), Mm(CT), Mm(CW), Mm(30))
            tf = txBox.text_frame
            setup_text_frame(tf)
            run = tf.paragraphs[0].add_run()
            set_run_font(run, f"Failed to render template '{spec.get('template', '?')}': {e}",
                         size_pt=10, color_rgb=CATHAY_RED)

    fixes = validate_and_fix(prs)
    prs.save(output_path)
    print(f"Deck saved: {len(prs.slides)} slides -> {output_path}")

    if fixes:
        print(f"  Auto-fixed {len(fixes)} issues")

    qc_report = {}
    if run_qc:
        try:
            qc_report = full_qc_pipeline(output_path)
        except Exception as e:
            print(f"  QC pipeline error: {e}")

    return {
        "path": output_path,
        "slides": len(prs.slides),
        "qc_report": qc_report,
        "fixes": fixes,
    }


__all__ = [
    "DataRegistry",
    "TEMPLATE_ROUTER",
    "render_spec",
    "build_deck_from_specs",
]
