"""
Cathay PPT Template — Merge & Anti-Corruption
===============================================
merge_slides(), reorder_slides(), clear_slide(), full_cleanup(), _clean_shape().

Usage:
    from merge import merge_slides, full_cleanup, clear_slide
"""

import os
import re
import copy
import io
import zipfile
import tempfile
import shutil

from pptx import Presentation
from pptx.oxml.ns import qn

from constants import TEMPLATE


# ============================================================================
# 1. ANTI-CORRUPTION DEFENSE
# ============================================================================

def _clean_shape(shape):
    """Strip <p:style> XML from shape to prevent theme corruption.
    Call immediately after creating any shape."""
    sp = shape._element
    for pstyle in sp.findall('.//' + qn('p:style')):
        pstyle.getparent().remove(pstyle)


def full_cleanup(pptx_path):
    """Post-save nuclear cleanup: strip ALL <p:style> and theme shadows from PPTX zip.
    Prevents PowerPoint theme corruption caused by python-pptx connector artifacts."""
    PSTYLE_RE = re.compile(r'<p:style>.*?</p:style>', re.DOTALL)

    tmp = tempfile.mktemp(suffix='.pptx')
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml') or item.filename.endswith('.rels'):
                    text = data.decode('utf-8')
                    text = PSTYLE_RE.sub('', text)
                    data = text.encode('utf-8')
                zout.writestr(item, data)
    shutil.move(tmp, pptx_path)


# ============================================================================
# 2. SLIDE MERGE
# ============================================================================

def merge_slides(slide_files, output_path, template_path=None, slide_order=None):
    """Merge multiple single-slide files into one deck with image rId remapping.

    Args:
        slide_files: dict {slide_num: path} or list of paths
        output_path: output file path
        template_path: template file path (defaults to TEMPLATE)
        slide_order: list of slide_nums for ordering

    Returns:
        int: number of slides in merged deck
    """
    template_path = template_path or TEMPLATE
    master = Presentation(template_path)

    # Clear template slides
    while len(master.slides) > 0:
        rId = master.slides._sldIdLst[0].rId
        master.part.drop_rel(rId)
        del master.slides._sldIdLst[0]

    if isinstance(slide_files, list):
        slide_files = {i + 1: p for i, p in enumerate(slide_files)}

    if slide_order is None:
        slide_order = sorted(slide_files.keys())

    for src_num in slide_order:
        if src_num not in slide_files:
            continue
        src_path = slide_files[src_num]
        if not os.path.exists(src_path):
            continue

        src_prs = Presentation(src_path)
        src_slide = src_prs.slides[0]

        # Match layout
        layout_name = src_slide.slide_layout.name
        target_layout = master.slide_layouts[4]
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                target_layout = layout
                break

        new_slide = master.slides.add_slide(target_layout)

        # Collect image blobs from source
        img_map = {}
        for rel in src_slide.part.rels.values():
            if "image" in str(rel.reltype):
                try:
                    img_map[rel.rId] = rel.target_part.blob
                except Exception:
                    pass

        # Register images in new slide, build rId mapping
        rId_remap = {}
        for old_rId, blob in img_map.items():
            image_part, new_rId = new_slide.part.get_or_add_image_part(
                io.BytesIO(blob))
            rId_remap[old_rId] = new_rId

        # Copy shapes with remapped image references
        for shape in src_slide.shapes:
            el = copy.deepcopy(shape._element)
            for blip in el.findall('.//' + qn('a:blip')):
                old_rId = blip.get(qn('r:embed'))
                if old_rId in rId_remap:
                    blip.set(qn('r:embed'), rId_remap[old_rId])
            new_slide.shapes._spTree.append(el)

    master.save(output_path)
    return len(master.slides)


# ============================================================================
# 3. SLIDE REORDER
# ============================================================================

def reorder_slides(prs, new_order_1based):
    """Reorder slides in a presentation.

    new_order_1based: list of 1-based slide indices in desired order.
    Example: [1, 2, 5, 3, 4] moves slide 5 to position 3.
    """
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    if len(new_order_1based) != len(ids):
        raise ValueError(
            f"new_order has {len(new_order_1based)} items but presentation has {len(ids)} slides")
    reordered = [ids[i - 1] for i in new_order_1based]
    for el in ids:
        sldIdLst.remove(el)
    for el in reordered:
        sldIdLst.append(el)


# ============================================================================
# 4. CLEAR SLIDE
# ============================================================================

def clear_slide(slide):
    """Remove all shapes from a slide (for rebuilding content on an existing slide).
    Preserves the slide's layout/master relationship."""
    for shp in list(slide.shapes):
        shp._element.getparent().remove(shp._element)


__all__ = [
    "_clean_shape",
    "full_cleanup",
    "merge_slides",
    "reorder_slides",
    "clear_slide",
]
