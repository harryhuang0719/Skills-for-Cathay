"""
Cathay PPT Template — Validation & QC
=======================================
validate_and_fix(), save_with_validation(), validate_no_overlap(),
validate_text_fit(), qc_presentation(), export_to_pdf().

Usage:
    from validation import validate_and_fix, save_with_validation, qc_presentation
"""

import os
import subprocess

from pptx import Presentation
from pptx.util import Pt

from constants import (
    SOURCE_Y_MM, SOURCE_BOX_HEIGHT_MM,
    MIN_SOURCE_FONT_PT, MIN_BODY_FONT_PT,
)
from text_layout import calc_textframe_height


# ============================================================================
# 1. PRE-SAVE VALIDATION + AUTO-FIX
# ============================================================================

def validate_and_fix(prs):
    """Pre-save validation: check every shape for overflow, auto-reduce font.

    Returns:
        list of fix descriptions
    """
    fixes = []

    for slide in prs.slides:
        for shape in slide.shapes:
            top_mm = shape.top / 36000
            height_mm = shape.height / 36000
            width_mm = shape.width / 36000
            bottom_mm = top_mm + height_mm

            if width_mm < 0.5 or height_mm < 0.5:
                continue

            # Check if it's a footer element
            is_footer = False
            if shape.has_text_frame:
                txt = shape.text_frame.text.lower()
                if 'source:' in txt or (len(txt) < 10 and '/' in txt):
                    is_footer = True

            # Cap shapes that exceed content zone (181mm)
            if not is_footer and bottom_mm > 181:
                new_h = 181 - top_mm
                if new_h >= 5:
                    shape.height = int(new_h * 36000)
                    fixes.append(f"CAP: {shape.name} bottom {bottom_mm:.0f}->181mm")

            # Check text overflow within textbox
            if shape.has_text_frame and height_mm >= 5:
                est_h = calc_textframe_height(shape.text_frame, width_mm)
                if est_h > height_mm * 1.1:
                    for target in [9, 8.5, 8, 7.5, 7]:
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                if r.font.size and r.font.size / 12700 > target:
                                    r.font.size = Pt(target)
                        new_est = calc_textframe_height(shape.text_frame, width_mm)
                        if new_est <= height_mm:
                            fixes.append(f"FONT: {shape.name} reduced to {target}pt")
                            break

    return fixes


def save_with_validation(prs, path):
    """Save with auto validation + fix + anti-corruption cleanup."""
    fixes = validate_and_fix(prs)
    if fixes:
        print(f"Auto-fixed {len(fixes)} issues before save:")
        for f in fixes[:10]:
            print(f"  {f}")
    prs.save(path)
    # Import here to avoid circular dependency
    from merge import full_cleanup
    full_cleanup(path)
    return fixes


# ============================================================================
# 2. OVERLAP DETECTION
# ============================================================================

def validate_no_overlap(pptx_path):
    """Check all slides for overlapping shapes. Returns list of issues."""
    prs_check = Presentation(pptx_path)
    issues = []
    for slide_idx, slide in enumerate(prs_check.slides, 1):
        shapes = []
        for sh in slide.shapes:
            l = sh.left / 914400
            t = sh.top / 914400
            r = l + sh.width / 914400
            b = t + sh.height / 914400
            shapes.append((sh.name, l, t, r, b))
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                n1, l1, t1, r1, b1 = shapes[i]
                n2, l2, t2, r2, b2 = shapes[j]
                if l1 < r2 and r1 > l2 and t1 < b2 and b1 > t2:
                    inside = ((l2 >= l1 and r2 <= r1 and t2 >= t1 and b2 <= b1) or
                              (l1 >= l2 and r1 <= r2 and t1 >= t2 and b1 <= b2))
                    if not inside:
                        issues.append(f"Slide {slide_idx}: '{n1}' overlaps '{n2}'")
    return issues


# ============================================================================
# 3. TEXT FIT ESTIMATION
# ============================================================================

def validate_text_fit(pptx_path):
    """Estimate whether text fits within each textbox. Returns warnings."""
    prs_check = Presentation(pptx_path)
    warnings = []
    for slide_idx, slide in enumerate(prs_check.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            total_text = "".join(p.text for p in tf.paragraphs)
            if not total_text.strip():
                continue
            box_w_mm = shape.width / 36000
            box_h_mm = shape.height / 36000
            has_chinese = any('\u4e00' <= c <= '\u9fff' for c in total_text)
            chars_per_mm = 2.5 if has_chinese else 3.5
            usable_w = box_w_mm - 4
            chars_per_line = max(usable_w * chars_per_mm, 1)
            num_paragraphs = len([p for p in tf.paragraphs if p.text.strip()])
            total_chars = len(total_text)
            est_lines = (total_chars / chars_per_line) + num_paragraphs * 0.3
            line_height_mm = 2.8
            est_height = est_lines * line_height_mm + 4
            if est_height > box_h_mm * 1.15:
                overflow_pct = ((est_height - box_h_mm) / box_h_mm) * 100
                warnings.append(
                    f"Slide {slide_idx}: '{shape.name}' text may overflow "
                    f"(est {est_height:.0f}mm vs box {box_h_mm:.0f}mm, +{overflow_pct:.0f}%)")
    return warnings


# ============================================================================
# 4. QC & PDF EXPORT
# ============================================================================

def qc_presentation(pptx_path):
    """Run full QC: overlap check + text fit check. Returns (issues, pdf_path)."""
    issues = validate_no_overlap(pptx_path)
    fit_warnings = validate_text_fit(pptx_path)
    if issues:
        print(f"OVERLAP ISSUES ({len(issues)}):")
        for issue in issues:
            print(f"  - {issue}")
    if fit_warnings:
        print(f"TEXT FIT WARNINGS ({len(fit_warnings)}):")
        for w in fit_warnings:
            print(f"  - {w}")
    if not issues and not fit_warnings:
        print("No layout issues found.")
    pdf = export_to_pdf(pptx_path)
    if pdf:
        print(f"PDF exported: {pdf}")
    return issues + fit_warnings, pdf


def export_to_pdf(pptx_path, output_dir=None):
    """Convert PPTX to PDF via LibreOffice for visual QC."""
    output_dir = output_dir or os.path.dirname(pptx_path)
    try:
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf',
             '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=120)
        pdf_path = os.path.join(
            output_dir,
            os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf')
        return pdf_path if os.path.exists(pdf_path) else None
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return None


__all__ = [
    "validate_and_fix",
    "save_with_validation",
    "validate_no_overlap",
    "validate_text_fit",
    "qc_presentation",
    "export_to_pdf",
]
