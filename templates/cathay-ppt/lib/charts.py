"""
Cathay PPT Template — Chart Helpers
=====================================
setup_chart_style(), safe_chart_insert(), insert_chart_image().

Usage:
    from charts import setup_chart_style, safe_chart_insert
"""

import os
import tempfile

from pptx.util import Mm

from constants import (
    CL, CT, CB, CW, CATHAY_COLORS, CONTENT_BOTTOM_MM,
)


def setup_chart_style():
    """Apply Cathay brand styling to matplotlib (call before creating charts)."""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import matplotlib.font_manager as fm

        available = {f.name for f in fm.fontManager.ttflist}
        font = 'Calibri' if 'Calibri' in available else 'Arial'
        plt.rcParams.update({
            'font.family': font,
            'font.size': 10,
            'axes.prop_cycle': plt.cycler(color=CATHAY_COLORS),
            'axes.edgecolor': '#808080',
            'axes.linewidth': 0.5,
            'grid.color': '#D9D9D9',
            'grid.linewidth': 0.5,
            'figure.facecolor': 'white',
            'axes.facecolor': 'white',
        })
    except ImportError:
        pass


def safe_chart_insert(slide, image_path, x_mm=None, y_mm=None, w_mm=200):
    """Insert chart PNG with width-only sizing, return actual bottom Y (mm).

    Reads actual PNG pixel dimensions, computes rendered height preserving
    aspect ratio. Auto-scales if chart would exceed content zone.

    Returns:
        bottom_y_mm (float): Y coordinate where the chart ends.
    """
    x_mm = x_mm or CL
    y_mm = y_mm or CT

    try:
        from PIL import Image
        with Image.open(image_path) as img:
            px_w, px_h = img.size
    except ImportError:
        px_w, px_h = 1600, 900

    aspect = px_h / px_w
    rendered_h_mm = w_mm * aspect

    _safe_bottom = CB - 3
    bottom_y = y_mm + rendered_h_mm
    if bottom_y > _safe_bottom:
        max_h = _safe_bottom - y_mm
        w_mm = max_h / aspect
        rendered_h_mm = max_h
        bottom_y = y_mm + rendered_h_mm

    slide.shapes.add_picture(image_path, Mm(x_mm), Mm(y_mm), Mm(w_mm))
    return bottom_y


def insert_chart_image(slide, image_path, x_mm=None, y_mm=None, w_mm=None):
    """Insert a chart image (width-only, preserves aspect ratio).
    DEPRECATED: prefer safe_chart_insert().
    """
    x_mm = x_mm or CL
    y_mm = y_mm or CT
    w_mm = w_mm or 200
    slide.shapes.add_picture(image_path, Mm(x_mm), Mm(y_mm), Mm(w_mm))


def cathay_bar_chart(categories, values, title, output_path=None, ylabel="", figsize=(8, 4.5)):
    """Cathay brand bar chart. Returns path to saved PNG."""
    setup_chart_style()
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=figsize)
    bars = ax.bar(categories, values, color=CATHAY_COLORS[:len(values)], width=0.6)
    ax.set_title(title, fontsize=12, fontweight='bold', color='#800000', loc='left')
    ax.set_ylabel(ylabel, fontsize=9, color='#808080')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.3)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(values)*0.02,
                f'{val:,.0f}' if isinstance(val, (int, float)) else str(val),
                ha='center', va='bottom', fontsize=9, color='#404040')
    plt.tight_layout()
    output_path = output_path or os.path.join(tempfile.gettempdir(), 'cathay_bar.png')
    fig.savefig(output_path, dpi=200, bbox_inches='tight')
    plt.close(fig)
    return output_path


def cathay_line_chart(x_labels, series_dict, title, output_path=None, ylabel="", figsize=(8, 4.5)):
    """Cathay brand line chart. series_dict: {"Series A": [vals], "Series B": [vals]}."""
    setup_chart_style()
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=figsize)
    for i, (name, vals) in enumerate(series_dict.items()):
        ax.plot(x_labels, vals, marker='o', markersize=4, linewidth=2,
                color=CATHAY_COLORS[i % len(CATHAY_COLORS)], label=name)
    ax.set_title(title, fontsize=12, fontweight='bold', color='#800000', loc='left')
    ax.set_ylabel(ylabel, fontsize=9, color='#808080')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.legend(fontsize=9, frameon=False)
    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    output_path = output_path or os.path.join(tempfile.gettempdir(), 'cathay_line.png')
    fig.savefig(output_path, dpi=200, bbox_inches='tight')
    plt.close(fig)
    return output_path


def cathay_waterfall_chart(labels, values, title, output_path=None, figsize=(8, 4.5)):
    """Waterfall chart (common in PE for value bridge / EV walk)."""
    setup_chart_style()
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=figsize)
    cumulative = [0]
    for v in values[:-1]:
        cumulative.append(cumulative[-1] + v)
    colors = []
    for i, v in enumerate(values):
        if i == 0 or i == len(values) - 1:
            colors.append('#800000')
        elif v >= 0:
            colors.append('#C8A415')
        else:
            colors.append('#E60000')
    bottoms = cumulative[:-1] + [0]
    ax.bar(labels, [abs(v) for v in values],
           bottom=[max(0, b) if i < len(values)-1 else 0 for i, b in enumerate(bottoms)],
           color=colors, width=0.5)
    ax.set_title(title, fontsize=12, fontweight='bold', color='#800000', loc='left')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    plt.xticks(rotation=30, ha='right', fontsize=9)
    plt.tight_layout()
    output_path = output_path or os.path.join(tempfile.gettempdir(), 'cathay_waterfall.png')
    fig.savefig(output_path, dpi=200, bbox_inches='tight')
    plt.close(fig)
    return output_path


__all__ = [
    "setup_chart_style",
    "safe_chart_insert",
    "insert_chart_image",
    "cathay_bar_chart",
    "cathay_line_chart",
    "cathay_waterfall_chart",
]
